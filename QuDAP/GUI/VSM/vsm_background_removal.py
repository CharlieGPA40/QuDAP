from PyQt6.QtWidgets import (
    QTreeWidgetItem, QMessageBox, QProgressBar, QWidget, QHeaderView, QGroupBox, QVBoxLayout, QLabel,
    QHBoxLayout, QSizePolicy, QTableWidget, QLineEdit, QPushButton, QMenu, QScrollArea, QTreeWidget,
    QWidgetAction, QMainWindow, QCheckBox, QFileDialog, QApplication)
from PyQt6.QtGui import QFont, QDragEnterEvent, QDropEvent
from PyQt6.QtCore import QPoint, Qt
import csv
import pandas as pd
import numpy as np
import matplotlib

matplotlib.use('QtAgg')
from matplotlib.backends.backend_qt5agg import FigureCanvasQTAgg as FigureCanvas
from matplotlib.backends.backend_qt5agg import NavigationToolbar2QT as NavigationToolbar
from matplotlib.figure import Figure
import matplotlib.pyplot as plt
import traceback
import os
import glob
import lmfit
from scipy import interpolate

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
except ImportError as e:
    print(e)


class MplCanvas(FigureCanvas):
    def __init__(self, parent=None, width=8, height=10, dpi=100):
        self.fig = Figure(figsize=(width, height), dpi=dpi)
        self.fig.clear()
        self.ax = self.fig.add_subplot(111, projection='rectilinear')
        super(MplCanvas, self).__init__(self.fig)


class DragDropWidget(QWidget):
    def __init__(self, main_window, label_text="Drag and drop a folder here"):
        super().__init__()
        self.main_window = main_window
        self.label_text = label_text
        self.initUI()

    def initUI(self):
        try:
            with open("GUI/SHG/QButtonWidget.qss", "r") as file:
                self.Browse_Button_stylesheet = file.read()
        except:
            self.Browse_Button_stylesheet = """
                QPushButton {
                    background-color: #4CAF50;
                    color: white;
                    border: none;
                    padding: 8px 16px;
                    border-radius: 4px;
                    font-size: 14px;
                }
                QPushButton:hover {
                    background-color: #45a049;
                }
            """

        self.setAcceptDrops(True)
        self.setStyleSheet("background-color: #F5F6FA; border: none;")

        self.label = QLabel(self.label_text, self)
        self.label.setAlignment(Qt.AlignmentFlag.AlignCenter)
        self.label.setStyleSheet("color: #D4D5D9; font-weight: bold; font-size: 16px;")

        self.button = QPushButton("Browse", self)
        self.button.setStyleSheet(self.Browse_Button_stylesheet)
        self.button.clicked.connect(self.open_folder_dialog)

        main_layout = QVBoxLayout(self)
        main_layout.addWidget(self.label, 5)
        main_layout.addWidget(self.button, 3, alignment=Qt.AlignmentFlag.AlignRight)
        self.setLayout(main_layout)

    def dragEnterEvent(self, event: QDragEnterEvent):
        try:
            if event.mimeData().hasUrls():
                event.acceptProposedAction()
            else:
                event.ignore()
        except Exception as e:
            QMessageBox.warning(self, "Error", str(e))

    def dropEvent(self, event: QDropEvent):
        try:
            urls = event.mimeData().urls()
            if urls:
                paths = [url.toLocalFile() for url in urls]
                directories = [path for path in paths if os.path.isdir(path)]
                csv_files = [path for path in paths if os.path.isfile(path) and path.lower().endswith('.csv')]

                if directories:
                    # Handle folder drop
                    self.main_window.set_folder(directories[0], self.label_text)
                elif csv_files:
                    # Handle CSV files drop
                    self.main_window.set_csv_files(csv_files, self.label_text)
                else:
                    QMessageBox.warning(self, "Invalid Drop", "Please drop a folder or CSV files.")
        except Exception as e:
            QMessageBox.warning(self, "Error", str(e))

    def open_folder_dialog(self):
        # Give user choice between folder or files
        choice_dialog = QMessageBox()
        choice_dialog.setWindowTitle("Select Input Type")
        choice_dialog.setText("How would you like to select your data?")
        folder_button = choice_dialog.addButton("Select Folder", QMessageBox.ButtonRole.AcceptRole)
        files_button = choice_dialog.addButton("Select CSV Files", QMessageBox.ButtonRole.AcceptRole)
        cancel_button = choice_dialog.addButton("Cancel", QMessageBox.ButtonRole.RejectRole)

        choice_dialog.exec()

        if choice_dialog.clickedButton() == folder_button:
            folder_path = QFileDialog.getExistingDirectory(self, "Select Folder")
            if folder_path:
                self.main_window.set_folder(folder_path, self.label_text)
        elif choice_dialog.clickedButton() == files_button:
            file_paths, _ = QFileDialog.getOpenFileNames(self, "Select CSV Files", "", "CSV Files (*.csv)")
            if file_paths:
                self.main_window.set_csv_files(file_paths, self.label_text)


class VSM_Background_Removal(QMainWindow):
    def __init__(self, parent=None):
        super().__init__(parent)
        self.isInit = False
        self.data_folder = None
        self.background_folder = None
        self.data_csv_files = []
        self.background_csv_files = []
        self.save_folder = None
        self.is_processed = False
        self.processed_files_data = {}
        self.init_ui()

    def init_ui(self):
        try:
            if self.isInit:
                return

            self.isInit = True

            # Load stylesheets
            try:
                with open("GUI/QSS/QScrollbar.qss", "r") as file:
                    self.scrollbar_stylesheet = file.read()
            except:
                self.scrollbar_stylesheet = ""

            try:
                with open("GUI/QSS/QButtonWidget.qss", "r") as file:
                    self.Button_stylesheet = file.read()
            except:
                self.Button_stylesheet = """
                    QPushButton {
                        background-color: #4CAF50;
                        color: white;
                        border: none;
                        padding: 8px 16px;
                        border-radius: 4px;
                        font-size: 14px;
                    }
                    QPushButton:hover {
                        background-color: #45a049;
                    }
                """

            titlefont = QFont("Arial", 20)
            self.font = QFont("Arial", 12)
            self.setStyleSheet("background-color: white;")

            # Create scroll area
            self.scroll_area = QScrollArea()
            self.scroll_area.setWidgetResizable(True)
            self.scroll_area.setHorizontalScrollBarPolicy(Qt.ScrollBarPolicy.ScrollBarAlwaysOff)
            self.scroll_area.setStyleSheet(self.scrollbar_stylesheet)

            self.content_widget = QWidget()
            self.scroll_area.setWidget(self.content_widget)

            self.main_layout = QVBoxLayout(self.content_widget)
            self.main_layout.setAlignment(Qt.AlignmentFlag.AlignCenter)
            self.content_widget.setSizePolicy(QSizePolicy.Policy.Expanding, QSizePolicy.Policy.Expanding)

            # Title
            self.title_label = QLabel("VSM Background Removal")
            self.title_label.setFont(titlefont)
            self.title_label.setStyleSheet("QLabel{ background-color: white; }")

            # File Upload Section
            self.fileUpload_layout = QHBoxLayout()

            # Data folder selection
            self.data_group_box = QGroupBox("Data Folder")
            self.data_group_box.setStyleSheet("QGroupBox { max-width: 550px; max-height: 230px; }")
            self.data_layout = QVBoxLayout()
            self.data_drag_drop = DragDropWidget(self, "Drag and drop DATA folder or CSV files here")
            self.data_status_label = QLabel('Please Select Data Folder or CSV Files')
            self.data_status_label.setStyleSheet("""
                color: white; font-size: 12px; background-color: #f38d76; 
                border-radius: 5px; padding: 5px;
            """)
            self.data_layout.addWidget(self.data_drag_drop, 4)
            self.data_layout.addWidget(self.data_status_label, 1, alignment=Qt.AlignmentFlag.AlignCenter)
            self.data_group_box.setLayout(self.data_layout)

            # Background folder selection
            self.bg_group_box = QGroupBox("Background Folder")
            self.bg_group_box.setStyleSheet("QGroupBox { max-width: 550px; max-height: 230px; }")
            self.bg_layout = QVBoxLayout()
            self.bg_drag_drop = DragDropWidget(self, "Drag and drop BACKGROUND folder or CSV files here")
            self.bg_status_label = QLabel('Please Select Background Folder or CSV Files')
            self.bg_status_label.setStyleSheet("""
                color: white; font-size: 12px; background-color: #f38d76; 
                border-radius: 5px; padding: 5px;
            """)
            self.bg_layout.addWidget(self.bg_drag_drop, 4)
            self.bg_layout.addWidget(self.bg_status_label, 1, alignment=Qt.AlignmentFlag.AlignCenter)
            self.bg_group_box.setLayout(self.bg_layout)

            self.fileUpload_layout.addWidget(self.data_group_box, 1)
            self.fileUpload_layout.addWidget(self.bg_group_box, 1)

            self.fileupload_container = QWidget(self)
            self.fileupload_container.setLayout(self.fileUpload_layout)
            self.fileupload_container.setFixedSize(1150, 300)

            # File Tree View Section (initially visible to show data files)
            self.file_tree_group_box = QGroupBox("Files")
            self.file_tree_group_box.setStyleSheet("QGroupBox { max-width: 1150px; max-height: 200px; }")
            self.file_tree_layout = QVBoxLayout()

            try:
                with open("GUI/SHG/QTreeWidget.qss", "r") as file:
                    self.QTree_stylesheet = file.read()
            except:
                self.QTree_stylesheet = ""

            self.file_tree = QTreeWidget()
            self.file_tree.setHeaderLabels(["Data File", "Temperature (K)", "Background Match", "Processing Status"])
            self.file_tree.setStyleSheet(self.QTree_stylesheet)
            self.file_tree.itemSelectionChanged.connect(self.on_file_tree_selection_changed)
            self.file_tree.setSelectionMode(QTreeWidget.SelectionMode.SingleSelection)

            self.file_tree_layout.addWidget(self.file_tree)
            self.file_tree_group_box.setLayout(self.file_tree_layout)
            # self.file_tree_group_box.hide()  # Hide initially until files are loaded

            self.file_tree_container = QWidget(self)
            self.file_tree_container.setLayout(QVBoxLayout())
            self.file_tree_container.layout().addWidget(self.file_tree_group_box)
            self.file_tree_container.setFixedSize(1150, 200)

            # Figure Layout
            self.figure_layout = QHBoxLayout()

            # Data Canvas
            self.data_canvas_layout = QVBoxLayout()
            self.data_canvas = MplCanvas(self, width=6, height=5, dpi=100)
            self.data_toolbar = NavigationToolbar(self.data_canvas, self)
            self.data_toolbar.setStyleSheet("QWidget { border: None; }")
            self.data_canvas.ax.set_title("Data Hysteresis")
            self.data_canvas_layout.addWidget(self.data_toolbar, alignment=Qt.AlignmentFlag.AlignCenter)
            self.data_canvas_layout.addWidget(self.data_canvas)

            # Background Canvas
            self.bg_canvas_layout = QVBoxLayout()
            self.bg_canvas = MplCanvas(self, width=6, height=5, dpi=100)
            self.bg_toolbar = NavigationToolbar(self.bg_canvas, self)
            self.bg_toolbar.setStyleSheet("QWidget { border: None; }")
            self.bg_canvas.ax.set_title("Background Hysteresis")
            self.bg_canvas_layout.addWidget(self.bg_toolbar, alignment=Qt.AlignmentFlag.AlignCenter)
            self.bg_canvas_layout.addWidget(self.bg_canvas)

            # Removed Canvas
            self.removed_canvas_layout = QVBoxLayout()
            self.removed_canvas = MplCanvas(self, width=6, height=5, dpi=100)
            self.removed_toolbar = NavigationToolbar(self.removed_canvas, self)
            self.removed_toolbar.setStyleSheet("QWidget { border: None; }")
            self.removed_canvas.ax.set_title("Background Removed")
            self.removed_canvas_layout.addWidget(self.removed_toolbar, alignment=Qt.AlignmentFlag.AlignCenter)
            self.removed_canvas_layout.addWidget(self.removed_canvas)

            self.figure_layout.addLayout(self.data_canvas_layout)
            self.figure_layout.addLayout(self.bg_canvas_layout)
            self.figure_layout.addLayout(self.removed_canvas_layout)

            self.figure_container = QWidget(self)
            self.figure_container.setLayout(self.figure_layout)
            self.figure_container.setFixedSize(1150, 400)

            # Button Layout
            self.button_layout = QHBoxLayout()
            self.process_button = QPushButton("Process")
            self.process_button.setStyleSheet(self.Button_stylesheet)
            self.process_button.clicked.connect(self.start_processing)
            self.process_button.setFixedSize(200, 40)

            self.reset_button = QPushButton("Reset")
            self.reset_button.setStyleSheet(self.Button_stylesheet)
            self.reset_button.clicked.connect(self.reset_all)
            self.reset_button.setFixedSize(200, 40)

            self.button_layout.addStretch()
            self.button_layout.addWidget(self.process_button)
            self.button_layout.addWidget(self.reset_button)

            self.button_container = QWidget(self)
            self.button_container.setLayout(self.button_layout)
            self.button_container.setFixedSize(1150, 50)

            # Progress Bar
            self.progress_layout = QHBoxLayout()
            self.progress_label = QLabel("Processing data...")
            self.progress_label.setStyleSheet("font-size: 12px; color: #4b6172;")
            self.progress_bar = QProgressBar()
            self.progress_bar.setStyleSheet("""
                QProgressBar {
                    border: 2px solid grey;
                    border-radius: 5px;
                    text-align: center;
                    height: 25px;
                }
                QProgressBar::chunk {
                    background-color: #4CAF50;
                    width: 10px;
                }
            """)
            self.progress_bar.setFixedWidth(750)
            self.progress_layout.addWidget(self.progress_label)
            self.progress_layout.addWidget(self.progress_bar)
            self.progress_layout.addStretch()

            # Hide progress bar initially
            self.progress_label.hide()
            self.progress_bar.hide()

            # Add to main layout
            self.main_layout.addWidget(self.title_label, alignment=Qt.AlignmentFlag.AlignTop)
            self.main_layout.addWidget(self.fileupload_container)
            self.main_layout.addWidget(self.button_container)
            self.main_layout.addWidget(self.file_tree_container)
            self.main_layout.addLayout(self.progress_layout)
            self.main_layout.addWidget(self.figure_container)
            self.main_layout.addStretch(1)
            self.setCentralWidget(self.scroll_area)

        except Exception as e:
            print(f"Error in init_ui: {e}")
            import traceback
            traceback.print_exc()
            QMessageBox.warning(self, "Error", f"Failed to initialize UI: {str(e)}")

    def set_folder(self, folder_path, label_text):
        """Set the folder path based on which drag-drop widget was used"""
        if "DATA" in label_text.upper():
            self.data_folder = folder_path
            self.data_csv_files = []
            self.data_status_label.setText(f"Data Folder: {os.path.basename(folder_path)}")
            self.data_status_label.setStyleSheet("""
                color: #4b6172; font-size: 12px; background-color: #DfE7Ef; 
                border-radius: 5px; padding: 5px;
            """)
            # Load and display data files in tree
            self.load_data_files_to_tree()
        elif "BACKGROUND" in label_text.upper():
            self.background_folder = folder_path
            self.background_csv_files = []
            self.bg_status_label.setText(f"Background Folder: {os.path.basename(folder_path)}")
            self.bg_status_label.setStyleSheet("""
                color: #4b6172; font-size: 12px; background-color: #DfE7Ef; 
                border-radius: 5px; padding: 5px;
            """)
            # Update background match status
            self.update_background_match_status()

    def set_csv_files(self, file_paths, label_text):
        """Set individual CSV files based on which drag-drop widget was used"""
        if "DATA" in label_text.upper():
            self.data_csv_files = file_paths
            self.data_folder = None
            file_count = len(file_paths)
            self.data_status_label.setText(f"Data: {file_count} CSV file{'s' if file_count > 1 else ''} selected")
            self.data_status_label.setStyleSheet("""
                color: #4b6172; font-size: 12px; background-color: #DfE7Ef; 
                border-radius: 5px; padding: 5px;
            """)
            # Load and display data files in tree
            self.load_data_files_to_tree()
        elif "BACKGROUND" in label_text.upper():
            self.background_csv_files = file_paths
            self.background_folder = None
            file_count = len(file_paths)
            self.bg_status_label.setText(f"Background: {file_count} CSV file{'s' if file_count > 1 else ''} selected")
            self.bg_status_label.setStyleSheet("""
                color: #4b6172; font-size: 12px; background-color: #DfE7Ef; 
                border-radius: 5px; padding: 5px;
            """)
            # Update background match status
            self.update_background_match_status()

    def load_data_files_to_tree(self):
        """Load data files into tree view and sort by temperature"""
        self.file_tree.clear()

        # Get data files
        if self.data_folder:
            data_csv_list_path = self.data_folder
            if not os.path.exists(data_csv_list_path):
                data_csv_list_path = self.data_folder + '/'
            data_csv_list = [f for f in os.listdir(data_csv_list_path) if f.endswith('.csv') and f != '.DS_Store']
        elif self.data_csv_files:
            data_csv_list = [os.path.basename(f) for f in self.data_csv_files]
        else:
            return

        # Create list of (filename, temperature) tuples for sorting
        file_temp_list = []
        for filename in data_csv_list:
            try:
                temp_str = filename.split('K')[0]
                temp_float = float(temp_str.replace('_', '.'))
                file_temp_list.append((filename, temp_float))
            except:
                file_temp_list.append((filename, 999999))  # Put unparseable files at the end

        # Sort by temperature
        file_temp_list.sort(key=lambda x: x[1])

        # Add items to tree
        for filename, temp in file_temp_list:
            item = QTreeWidgetItem(self.file_tree, [
                filename,
                str(temp) if temp != 999999 else "N/A",
                "Checking...",
                "Not Processed"
            ])
            temp_key = filename.split('K')[0]
            item.setData(0, Qt.ItemDataRole.UserRole, temp_key)

        # Show tree view
        self.file_tree_group_box.show()

        # Resize columns
        self.file_tree.resizeColumnToContents(0)
        self.file_tree.resizeColumnToContents(1)
        self.file_tree.resizeColumnToContents(2)
        self.file_tree.resizeColumnToContents(3)

        # Check background matches
        self.update_background_match_status()

    def update_background_match_status(self):
        """Update the background match status column"""
        if self.file_tree.topLevelItemCount() == 0:
            return

        # Get background files
        if self.background_folder:
            bg_csv_list_path = self.background_folder
            if not os.path.exists(bg_csv_list_path):
                bg_csv_list_path = self.background_folder + '/'
            bg_csv_list = [f for f in os.listdir(bg_csv_list_path) if f.endswith('.csv') and f != '.DS_Store']
        elif self.background_csv_files:
            bg_csv_list = [os.path.basename(f) for f in self.background_csv_files]
        else:
            # No background files selected yet
            for i in range(self.file_tree.topLevelItemCount()):
                item = self.file_tree.topLevelItem(i)
                item.setText(2, "No Background Selected")
            return

        # Check each data file for matching background
        for i in range(self.file_tree.topLevelItemCount()):
            item = self.file_tree.topLevelItem(i)
            data_filename = item.text(0)

            if data_filename in bg_csv_list:
                item.setText(2, "✓ Match Found")
            else:
                item.setText(2, "✗ No Match")

        # Resize column
        self.file_tree.resizeColumnToContents(2)

    def start_processing(self):
        """Start the background removal processing"""
        if not self.data_folder and not self.data_csv_files:
            QMessageBox.warning(self, "Missing Data", "Please select a data folder or CSV files.")
            return

        if not self.background_folder and not self.background_csv_files:
            QMessageBox.warning(self, "Missing Background", "Please select a background folder or CSV files.")
            return

        self.save_folder = QFileDialog.getExistingDirectory(self, "Select Save Location")
        if not self.save_folder:
            return

        self.progress_label.show()
        self.progress_bar.show()
        self.progress_bar.setValue(0)
        self.process_button.setEnabled(False)

        try:
            self.process_background_removal()
        except Exception as e:
            tb_str = traceback.format_exc()
            QMessageBox.warning(self, "Processing Error",
                                f"Error during processing:\n{str(e)}\n\n{str(tb_str)}")
        finally:
            self.progress_label.hide()
            self.progress_bar.hide()
            self.process_button.setEnabled(True)

    def splitting(self, input_df):
        """Split hysteresis data into upper and lower branches"""
        index = input_df.iloc[:, 0].idxmin(axis=0)
        list1 = input_df.iloc[:index].dropna()
        list2 = input_df.iloc[index:].dropna()

        list1_x = list1.iloc[:, 0].values
        list1_y = list1.iloc[:, 1].values
        list2_x = list2.iloc[:, 0].values
        list2_y = list2.iloc[:, 1].values
        return list1_x, list1_y, list2_x, list2_y

    def process_background_removal(self):
        """Main processing function"""
        self.progress_bar.setValue(5)
        QApplication.processEvents()

        self.BG_Removal_Data = self.save_folder + '/Background_Removal_Results'
        os.makedirs(self.BG_Removal_Data, exist_ok=True)

        self.Processed_Data = self.BG_Removal_Data + '/Processed_Data'
        os.makedirs(self.Processed_Data, exist_ok=True)

        self.bg_removal_RAW = self.BG_Removal_Data + '/BG_Removal_RAW'
        os.makedirs(self.bg_removal_RAW, exist_ok=True)

        self.BG_Removal_Processed_Graph = self.BG_Removal_Data + '/Processed_Graph'
        os.makedirs(self.BG_Removal_Processed_Graph, exist_ok=True)

        self.BG_Removal_ProcCal = self.BG_Removal_Data + '/Processed_Result'
        os.makedirs(self.BG_Removal_ProcCal, exist_ok=True)

        self.progress_bar.setValue(10)
        QApplication.processEvents()

        if self.data_folder:
            data_csv_list_path = self.data_folder + '/'
            if not os.path.exists(data_csv_list_path):
                data_csv_list_path = self.data_folder + '/'

            data_csv_list = [f for f in os.listdir(data_csv_list_path) if f.endswith('.csv') and f != '.DS_Store']
            data_files_dict = {f: data_csv_list_path + f for f in data_csv_list}
        else:
            data_csv_list = [os.path.basename(f) for f in self.data_csv_files]
            data_files_dict = {os.path.basename(f): f for f in self.data_csv_files}

        data_number_CSV = len(data_csv_list)
        if data_number_CSV == 0:
            QMessageBox.warning(self, "No Files", "No CSV files found in data input.")
            return

        if self.background_folder:
            bg_csv_list_path = self.background_folder + '/'
            if not os.path.exists(bg_csv_list_path):
                bg_csv_list_path = self.background_folder + '/'

            bg_csv_list = [f for f in os.listdir(bg_csv_list_path) if f.endswith('.csv') and f != '.DS_Store']
            bg_files_dict = {f: bg_csv_list_path + f for f in bg_csv_list}
        else:
            bg_csv_list = [os.path.basename(f) for f in self.background_csv_files]
            bg_files_dict = {os.path.basename(f): f for f in self.background_csv_files}

        self.progress_bar.setValue(15)
        QApplication.processEvents()

        for i in range(data_number_CSV):
            progress = 15 + int((i / data_number_CSV) * 70)
            self.progress_bar.setValue(progress)
            self.progress_label.setText(f"Processing file {i + 1}/{data_number_CSV}...")
            QApplication.processEvents()

            if data_csv_list[i] == '.DS_Store':
                continue

            j = 0
            NotMatch = True
            while NotMatch and j < len(bg_csv_list):
                if data_csv_list[i] == bg_csv_list[j]:
                    NotMatch = False
                else:
                    j += 1

            if NotMatch:
                print(f"Warning: No matching background file for {data_csv_list[i]}")
                continue

            Cur_Temp = data_csv_list[i].split('K')[0]

            try:
                temp_float = float(Cur_Temp.replace('_', '.'))
                cur_temp_numeric = temp_float
            except:
                cur_temp_numeric = float(Cur_Temp)

            try:
                data_file_path = data_files_dict[data_csv_list[i]]
                bg_file_path = bg_files_dict[bg_csv_list[j]]

                Data_Hyst = pd.read_csv(data_file_path, header=None, engine='c')
                bg_Hyst = pd.read_csv(bg_file_path, header=None, engine='c')

                x_data_top, y_data_top, x_data_bottom, y_data_bottom = self.splitting(Data_Hyst)
                x_bg_top, y_bg_top, x_bg_bottom, y_bg_bottom = self.splitting(bg_Hyst)

                top_interp_result = interpolate.interp1d(x_bg_top, y_bg_top, 'linear', fill_value="extrapolate")
                y_bg_new_top = top_interp_result(x_data_top)
                y_processed_top = y_data_top - y_bg_new_top

                bot_interp_result = interpolate.interp1d(x_bg_bottom, y_bg_bottom, 'linear', fill_value="extrapolate")
                y_bg_new_bot = bot_interp_result(x_data_bottom)
                y_processed_bot = y_data_bottom - y_bg_new_bot

                df_top = pd.DataFrame({'x': x_data_top, 'y': y_processed_top})
                z_scores = np.abs((df_top['y'] - df_top['y'].mean()) / df_top['y'].std())
                threshold = 3
                filtered_df = df_top[z_scores <= threshold]
                x_data_top = filtered_df['x'].values
                y_processed_top = filtered_df['y'].values

                df_bot = pd.DataFrame({'x': x_data_bottom, 'y': y_processed_bot})
                z_scores = np.abs((df_bot['y'] - df_bot['y'].mean()) / df_bot['y'].std())
                filtered_df = df_bot[z_scores <= threshold]
                x_data_bottom = filtered_df['x'].values
                y_processed_bot = filtered_df['y'].values

                if i == 0:
                    self.plot_comparison(Data_Hyst, bg_Hyst,
                                         x_data_top, y_processed_top,
                                         x_data_bottom, y_processed_bot, Cur_Temp)

                x_all = np.concatenate([x_data_top, x_data_bottom])
                y_all = np.concatenate([y_processed_top, y_processed_bot])

                output_df = pd.DataFrame({
                    'Magnetic Field (Oe)': x_all,
                    'Moment (emu)': y_all
                })

                output_filename = f'{Cur_Temp}K.csv'
                output_df.to_csv(self.Processed_Data + '/' + output_filename, index=False, header=True)

                bg_removal_df = pd.DataFrame({'Field': x_all, 'Moment': y_all})
                bg_removal_df.to_csv(
                    self.bg_removal_RAW + f'/{Cur_Temp}K_background_removal.csv',
                    index=False, header=False)

                self.save_individual_plot(x_data_top, y_processed_top,
                                          x_data_bottom, y_processed_bot, Cur_Temp)

                x_data = Data_Hyst.iloc[:, 0].values
                y_data = Data_Hyst.iloc[:, 1].values
                x_bg = bg_Hyst.iloc[:, 0].values
                y_bg = bg_Hyst.iloc[:, 1].values

                self.processed_files_data[Cur_Temp] = {
                    'data_x': x_data,
                    'data_y': y_data,
                    'bg_x': x_bg,
                    'bg_y': y_bg,
                    'removed_top_x': x_data_top,
                    'removed_top_y': y_processed_top,
                    'removed_bot_x': x_data_bottom,
                    'removed_bot_y': y_processed_bot,
                    'filename': data_csv_list[i]
                }

                # Update existing tree item instead of creating new one
                for idx in range(self.file_tree.topLevelItemCount()):
                    tree_item = self.file_tree.topLevelItem(idx)
                    if tree_item.text(0) == data_csv_list[i]:
                        tree_item.setText(3, "✓ Processed")
                        break

            except Exception as e:
                print(f"Error processing {data_csv_list[i]}: {e}")
                traceback.print_exc()

                # Update existing tree item to show failure
                for idx in range(self.file_tree.topLevelItemCount()):
                    tree_item = self.file_tree.topLevelItem(idx)
                    if tree_item.text(0) == data_csv_list[i]:
                        tree_item.setText(3, "✗ Failed")
                        break
                continue

        self.progress_bar.setValue(90)
        self.progress_label.setText("Generating plots...")
        QApplication.processEvents()

        self.to_ppt()

        self.progress_bar.setValue(100)
        self.progress_label.setText("Processing complete!")
        QApplication.processEvents()

        self.is_processed = True

        # Resize columns after processing
        self.file_tree.resizeColumnToContents(0)
        self.file_tree.resizeColumnToContents(1)
        self.file_tree.resizeColumnToContents(2)
        self.file_tree.resizeColumnToContents(3)

        QMessageBox.information(self, "Processing Complete",
                                f"Successfully processed {data_number_CSV} files.\n\n" +
                                f"Results saved to:\n{self.BG_Removal_Data}\n\n" +
                                "Click on files in the tree view to visualize the results.")

    def plot_comparison(self, data_df, bg_df, x_removed_top, y_removed_top,
                        x_removed_bot, y_removed_bot, temp):
        """Plot data, background, and removed result side by side"""
        self.data_canvas.fig.clear()
        self.data_canvas.ax = self.data_canvas.fig.add_subplot(111)

        self.bg_canvas.fig.clear()
        self.bg_canvas.ax = self.bg_canvas.fig.add_subplot(111)

        self.removed_canvas.fig.clear()
        self.removed_canvas.ax = self.removed_canvas.fig.add_subplot(111)

        x_data = data_df.iloc[:, 0].values
        y_data = data_df.iloc[:, 1].values
        self.data_canvas.ax.scatter(x_data, y_data, s=1, alpha=0.6, color='blue', label='Data')
        self.data_canvas.ax.set_xlabel('Magnetic Field (Oe)', fontsize=11)
        self.data_canvas.ax.set_ylabel('Moment (emu)', fontsize=11)
        self.data_canvas.ax.set_title(f'{temp}K Data Hysteresis', fontsize=12)
        self.data_canvas.ax.grid(True, alpha=0.3)
        self.data_canvas.ax.tick_params(labelsize=9)
        handles, labels = self.data_canvas.ax.get_legend_handles_labels()
        if handles:
            self.data_canvas.ax.legend(fontsize=8, loc='best')
        self.data_canvas.fig.tight_layout()
        self.data_canvas.draw()

        x_bg = bg_df.iloc[:, 0].values
        y_bg = bg_df.iloc[:, 1].values
        self.bg_canvas.ax.scatter(x_bg, y_bg, s=1, alpha=0.6, color='red', label='Background')
        self.bg_canvas.ax.set_xlabel('Magnetic Field (Oe)', fontsize=11)
        self.bg_canvas.ax.set_ylabel('Moment (emu)', fontsize=11)
        self.bg_canvas.ax.set_title(f'{temp}K Background Hysteresis', fontsize=12)
        self.bg_canvas.ax.grid(True, alpha=0.3)
        self.bg_canvas.ax.tick_params(labelsize=9)
        handles, labels = self.bg_canvas.ax.get_legend_handles_labels()
        if handles:
            self.bg_canvas.ax.legend(fontsize=8, loc='best')
        self.bg_canvas.fig.tight_layout()
        self.bg_canvas.draw()

        self.removed_canvas.ax.scatter(x_removed_top, y_removed_top, s=2, alpha=0.6,
                                       color='green', label='Upper Branch')
        self.removed_canvas.ax.scatter(x_removed_bot, y_removed_bot, s=2, alpha=0.6,
                                       color='orange', label='Lower Branch')
        self.removed_canvas.ax.set_xlabel('Magnetic Field (Oe)', fontsize=11)
        self.removed_canvas.ax.set_ylabel('Moment (emu)', fontsize=11)
        self.removed_canvas.ax.set_title(f'{temp}K Background Removed', fontsize=12)
        self.removed_canvas.ax.grid(True, alpha=0.3)
        self.removed_canvas.ax.tick_params(labelsize=9)
        handles, labels = self.removed_canvas.ax.get_legend_handles_labels()
        if handles:
            self.removed_canvas.ax.legend(fontsize=8, loc='best')
        self.removed_canvas.fig.tight_layout()
        self.removed_canvas.draw()

    def save_individual_plot(self, x_top, y_top, x_bot, y_bot, temp):
        """Save individual plot for each temperature"""
        try:
            plots_dir = self.BG_Removal_Processed_Graph + '/Individual_Plots'
            os.makedirs(plots_dir, exist_ok=True)

            fig, ax = plt.subplots(figsize=(8, 6))
            ax.scatter(x_top, y_top, s=2, alpha=0.6, color='cornflowerblue', label='Upper Branch')
            ax.scatter(x_bot, y_bot, s=2, alpha=0.6, color='darkorange', label='Lower Branch')
            ax.set_xlabel('Magnetic Field (Oe)', fontsize=14)
            ax.set_ylabel('Moment (emu)', fontsize=14)
            ax.set_title(f'{temp}K Background Removed Hysteresis', fontsize=14)
            ax.grid(True, alpha=0.3)
            ax.legend(fontsize=10)
            plt.tight_layout()
            plt.savefig(f'{plots_dir}/{temp}K_background_removed.png', dpi=150)
            plt.close()
        except Exception as e:
            print(f"Error saving plot for {temp}K: {e}")

    def on_file_tree_selection_changed(self):
        """Handle tree item selection change to display plots"""
        if not self.is_processed:
            return

        selected_items = self.file_tree.selectedItems()
        if not selected_items:
            return

        item = selected_items[0]
        temp_key = item.data(0, Qt.ItemDataRole.UserRole)

        if temp_key not in self.processed_files_data:
            return

        data = self.processed_files_data[temp_key]
        self.update_plots_from_data(data, temp_key)

    def update_plots_from_data(self, data, temp):
        """Update all three canvases with stored data"""
        self.data_canvas.fig.clear()
        self.data_canvas.ax = self.data_canvas.fig.add_subplot(111)

        self.bg_canvas.fig.clear()
        self.bg_canvas.ax = self.bg_canvas.fig.add_subplot(111)

        self.removed_canvas.fig.clear()
        self.removed_canvas.ax = self.removed_canvas.fig.add_subplot(111)

        self.data_canvas.ax.scatter(data['data_x'], data['data_y'],
                                    s=1, alpha=0.6, color='blue', label='Data')
        self.data_canvas.ax.set_xlabel('Magnetic Field (Oe)', fontsize=11)
        self.data_canvas.ax.set_ylabel('Moment (emu)', fontsize=11)
        self.data_canvas.ax.set_title(f'{temp}K Data Hysteresis', fontsize=12)
        self.data_canvas.ax.grid(True, alpha=0.3)
        self.data_canvas.ax.tick_params(labelsize=9)
        handles, labels = self.data_canvas.ax.get_legend_handles_labels()
        if handles:
            self.data_canvas.ax.legend(fontsize=8, loc='best')
        self.data_canvas.fig.tight_layout()
        self.data_canvas.draw()

        self.bg_canvas.ax.scatter(data['bg_x'], data['bg_y'],
                                  s=1, alpha=0.6, color='red', label='Background')
        self.bg_canvas.ax.set_xlabel('Magnetic Field (Oe)', fontsize=11)
        self.bg_canvas.ax.set_ylabel('Moment (emu)', fontsize=11)
        self.bg_canvas.ax.set_title(f'{temp}K Background Hysteresis', fontsize=12)
        self.bg_canvas.ax.grid(True, alpha=0.3)
        self.bg_canvas.ax.tick_params(labelsize=9)
        handles, labels = self.bg_canvas.ax.get_legend_handles_labels()
        if handles:
            self.bg_canvas.ax.legend(fontsize=8, loc='best')
        self.bg_canvas.fig.tight_layout()
        self.bg_canvas.draw()

        self.removed_canvas.ax.scatter(data['removed_top_x'], data['removed_top_y'],
                                       s=2, alpha=0.6, color='green', label='Upper Branch')
        self.removed_canvas.ax.scatter(data['removed_bot_x'], data['removed_bot_y'],
                                       s=2, alpha=0.6, color='orange', label='Lower Branch')
        self.removed_canvas.ax.set_xlabel('Magnetic Field (Oe)', fontsize=11)
        self.removed_canvas.ax.set_ylabel('Moment (emu)', fontsize=11)
        self.removed_canvas.ax.set_title(f'{temp}K Background Removed', fontsize=12)
        self.removed_canvas.ax.grid(True, alpha=0.3)
        self.removed_canvas.ax.tick_params(labelsize=9)
        handles, labels = self.removed_canvas.ax.get_legend_handles_labels()
        if handles:
            self.removed_canvas.ax.legend(fontsize=8, loc='best')
        self.removed_canvas.fig.tight_layout()
        self.removed_canvas.draw()

    def to_ppt(self):
        """Generate PowerPoint presentation with all processed plots"""
        try:
            prs = Presentation()
            prs.slide_width = Inches(13.33)
            prs.slide_height = Inches(7.5)

            blank_slide_layout = prs.slide_layouts[6]
            slide = prs.slides.add_slide(blank_slide_layout)

            text_frame = slide.shapes.add_textbox(Inches(5.67), Inches(0.06), Inches(2.03), Inches(0.57))
            text_frame = text_frame.text_frame
            p = text_frame.paragraphs[0]
            run = p.add_run()
            run.text = "Background Removal Results"
            font = run.font
            font.name = 'Calibri'
            font.size = Pt(28)

            plots_dir = self.BG_Removal_Processed_Graph + '/Individual_Plots'
            if os.path.exists(plots_dir):
                png_list = glob.glob(plots_dir + '/*.png')

                if png_list:
                    def extract_number(f):
                        try:
                            basename = os.path.basename(f)
                            return float(basename.split('K')[0].replace('_', '.'))
                        except:
                            return 0

                    png_list = sorted(png_list, key=extract_number)

                    initial_x = 0
                    initial_y = 0.89
                    image_width = 2.15
                    x_offset = 2.24
                    y_offset = 1.67
                    iteration = 0

                    for png_file in png_list:
                        if iteration == 6:
                            initial_x = 0
                            initial_y = initial_y + y_offset
                            iteration = 0

                        if iteration == 0 and initial_y > 0.89:
                            slide = prs.slides.add_slide(blank_slide_layout)
                            initial_x = 0
                            initial_y = 0.89

                        slide.shapes.add_picture(png_file, Inches(initial_x),
                                                 Inches(initial_y), Inches(image_width))
                        initial_x += x_offset
                        iteration += 1

            prs.save(self.BG_Removal_ProcCal + '/Background_Removal_Results.pptx')
            print("PowerPoint presentation saved successfully")
        except Exception as e:
            print(f"Error generating PowerPoint: {e}")
            traceback.print_exc()

    def reset_all(self):
        """Reset all data and clear the interface"""
        reply = QMessageBox.question(self, 'Reset Confirmation',
                                     'Are you sure you want to reset?',
                                     QMessageBox.StandardButton.Yes | QMessageBox.StandardButton.No,
                                     QMessageBox.StandardButton.No)

        if reply == QMessageBox.StandardButton.Yes:
            self.data_canvas.ax.clear()
            self.data_canvas.ax.set_title("Data Hysteresis")
            self.data_canvas.draw()

            self.bg_canvas.ax.clear()
            self.bg_canvas.ax.set_title("Background Hysteresis")
            self.bg_canvas.draw()

            self.removed_canvas.ax.clear()
            self.removed_canvas.ax.set_title("Background Removed")
            self.removed_canvas.draw()

            self.file_tree.clear()
            # self.file_tree_group_box.hide()
            self.processed_files_data = {}

            self.data_folder = None
            self.background_folder = None
            self.data_csv_files = []
            self.background_csv_files = []
            self.save_folder = None
            self.is_processed = False

            self.data_status_label.setText('Please Select Data Folder or CSV Files')
            self.data_status_label.setStyleSheet("""
                color: white; font-size: 12px; background-color: #f38d76; 
                border-radius: 5px; padding: 5px;
            """)

            self.bg_status_label.setText('Please Select Background Folder or CSV Files')
            self.bg_status_label.setStyleSheet("""
                color: white; font-size: 12px; background-color: #f38d76; 
                border-radius: 5px; padding: 5px;
            """)


if __name__ == "__main__":
    import sys
    from PyQt6.QtWidgets import QApplication

    app = QApplication(sys.argv)
    window = VSM_Background_Removal()
    window.show()
    sys.exit(app.exec())