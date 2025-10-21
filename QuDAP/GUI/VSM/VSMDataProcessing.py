from PyQt6.QtWidgets import (
    QTreeWidgetItem, QMessageBox, QProgressBar, QWidget, QHeaderView, QGroupBox, QVBoxLayout, QLabel,
    QHBoxLayout, QSizePolicy, QTableWidget, QLineEdit, QPushButton, QMenu, QScrollArea, QTreeWidget,
    QWidgetAction, QMainWindow, QCheckBox, QFileDialog)
from PyQt6.QtGui import QFont, QDragEnterEvent, QDropEvent
from PyQt6.QtCore import QPoint, Qt
import csv
import pandas as pd
import numpy as np
import matplotlib

matplotlib.use('QtAgg')
from matplotlib.backends.backend_qt5agg import FigureCanvasQTAgg as FigureCanvas
from matplotlib.backends.backend_qt5agg import NavigationToolbar2QT as NavigationToolbar
from matplotlib.figure import Figure
import matplotlib.pyplot as plt
import traceback
import os
import glob
import lmfit

import platform

system = platform.system()

try:
    from QuDAP.GUI.VSM.qd import *
    from QuDAP.misc.vsm_magnetic_contribution import (
        extract_pm_dm_slope,
        remove_pm_dm_contribution,
        visualize_pm_dm_extraction
    )
    from QuDAP.misc.vsm_tanh_fit import (
        process_vsm_branches,
        tanh_model,
        correct_vertical_offset,
        extract_saturation_magnetization,
        extract_coercivity_method1_fit,
        extract_coercivity_method2_data
    )
except ImportError:
    try:
        from GUI.VSM.qd import *
        from misc.vsm_fitting_method import (
        extract_pm_dm_slope,
        remove_pm_dm_contribution,
        visualize_pm_dm_extraction
        )
        from misc.vsm_tanh_fit import (
            process_vsm_branches,
            tanh_model,
            correct_vertical_offset,
            extract_saturation_magnetization,
            extract_coercivity_method1_fit,
            extract_coercivity_method2_data
        )
    except ImportError:
        print("Warning: qd module not found")

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
except ImportError as e:
    print(e)


class MplCanvas(FigureCanvas):
    def __init__(self, parent=None, width=8, height=10, dpi=1000):
        self.fig = Figure(figsize=(width, height), dpi=dpi)
        self.fig.clear()
        self.ax = self.fig.add_subplot(111, projection='rectilinear')
        super(MplCanvas, self).__init__(self.fig)


class DragDropWidget(QWidget):
    def __init__(self, main_window):
        super().__init__()
        self.main_window = main_window
        self.initUI()

    def initUI(self):
        try:
            with open("GUI/SHG/QButtonWidget.qss", "r") as file:
                self.Browse_Button_stylesheet = file.read()
        except:
            self.Browse_Button_stylesheet = ""

        self.setAcceptDrops(True)
        self.previous_folder_path = None
        self.setStyleSheet("background-color: #F5F6FA; border: none;")

        self.label = QLabel("Drag and drop a csv file or folder here", self)
        self.label.setAlignment(Qt.AlignmentFlag.AlignCenter)
        self.label.setStyleSheet("color: #D4D5D9; font-weight: bold; font-size: 16px;")

        self.button = QPushButton("Browse", self)
        self.button.setStyleSheet(self.Browse_Button_stylesheet)
        self.button.clicked.connect(self.open_folder_dialog)
        main_layout = QVBoxLayout(self)
        main_layout.addWidget(self.label, 5)
        main_layout.addWidget(self.button, 3, alignment=Qt.AlignmentFlag.AlignRight)
        self.setLayout(main_layout)

    def dragEnterEvent(self, event: QDragEnterEvent):
        try:
            if event.mimeData().hasUrls():
                event.acceptProposedAction()
            else:
                event.ignore()
        except Exception as e:
            QMessageBox.warning(self, "Error", str(e))
            return

    def dropEvent(self, event: QDropEvent):
        try:
            urls = event.mimeData().urls()
            if urls:
                paths = [url.toLocalFile() for url in urls]
                directories = [path for path in paths if os.path.isdir(path)]
                dat_files = [path for path in paths if os.path.isfile(path) and path.lower().endswith('.csv')]

                if directories:
                    for directory in directories:
                        self.main_window.display_files(directory + '/')

                if dat_files:
                    self.main_window.display_multiple_files(dat_files)

                if not directories and not dat_files:
                    QMessageBox.warning(self, "Invalid File", "Please drop a .csv file or a folder.")
        except Exception as e:
            QMessageBox.warning(self, "Error", str(e))
            return

    def open_folder_dialog(self):
        folder_path = QFileDialog.getExistingDirectory(self, "Select Folder")
        if folder_path:
            self.main_window.display_files(folder_path + '/')

    def reset(self):
        self.previous_folder_path = None

class VSM_Data_Processing(QMainWindow):
    def __init__(self, parent=None):
        super().__init__(parent)
        try:
            self.isInit = False
            self.file_in_list = []
            self.current_file_data = None
            self.area_df = None
            self.Ms_df = None
            self.Coercivity_df = None
            self.eb_df = None
            self.folder_selected = None
            self.ProcCal = None
            self.ProcessedRAW = None
            self.is_processed = False
            self.init_ui()
        except Exception as e:
            print(f"Initialization error: {e}")
            QMessageBox.warning(self, "Error", str(e))
            return

    def init_ui(self):
        try:
            if self.isInit == False:
                print("Starting UI initialization...")
                self.isInit = True

                # Load stylesheets with fallback
                try:
                    with open("GUI/VSM/QButtonWidget.qss", "r") as file:
                        self.Browse_Button_stylesheet = file.read()
                except:
                    self.Browse_Button_stylesheet = """
                        QPushButton {
                            background-color: #4CAF50;
                            color: white;
                            border: none;
                            padding: 8px 16px;
                            border-radius: 4px;
                            font-size: 14px;
                        }
                        QPushButton:hover {
                            background-color: #45a049;
                        }
                    """

                try:
                    with open("GUI/QSS/QScrollbar.qss", "r") as file:
                        self.scrollbar_stylesheet = file.read()
                except:
                    self.scrollbar_stylesheet = ""

                try:
                    with open("GUI/QSS/QButtonWidget.qss", "r") as file:
                        self.Button_stylesheet = file.read()
                except:
                    self.Button_stylesheet = self.Browse_Button_stylesheet

                try:
                    with open("GUI/SHG/QTreeWidget.qss", "r") as file:
                        self.QTree_stylesheet = file.read()
                except:
                    self.QTree_stylesheet = ""

                titlefont = QFont("Arial", 20)
                self.font = QFont("Arial", 12)
                self.setStyleSheet("background-color: white;")

                # Create scroll area
                self.scroll_area = QScrollArea()
                self.scroll_area.setWidgetResizable(True)
                self.scroll_area.setHorizontalScrollBarPolicy(Qt.ScrollBarPolicy.ScrollBarAlwaysOff)
                self.scroll_area.setStyleSheet(self.scrollbar_stylesheet)

                self.content_widget = QWidget()
                self.scroll_area.setWidget(self.content_widget)

                self.VSM_data_extraction_main_layout = QVBoxLayout(self.content_widget)
                self.VSM_data_extraction_main_layout.setAlignment(Qt.AlignmentFlag.AlignCenter)
                self.content_widget.setSizePolicy(QSizePolicy.Policy.Expanding, QSizePolicy.Policy.Expanding)

                # Title
                self.VSM_QD_EXTRACT_label = QLabel("VSM Data Processing")
                self.VSM_QD_EXTRACT_label.setFont(titlefont)
                self.VSM_QD_EXTRACT_label.setStyleSheet("QLabel{ background-color: white; }")

                # File Upload Section
                self.fileUpload_layout = QHBoxLayout()
                self.drag_drop_layout = QVBoxLayout()
                self.file_selection_group_box = QGroupBox("Upload Directory")
                self.file_view_group_box = QGroupBox("View Files")

                self.file_selection_group_box.setStyleSheet("""
                    QGroupBox { max-width: 600px; max-height: 230px; }
                """)
                self.file_view_group_box.setStyleSheet("""
                    QGroupBox { max-width: 650px; max-height: 230px; }
                """)

                self.file_selection_display_label = QLabel('Please Upload Files or Directory')
                self.file_selection_display_label.setStyleSheet("""
                    color: white; font-size: 12px; background-color: #f38d76; 
                    border-radius: 5px; padding: 5px;
                """)

                self.drag_drop_widget = DragDropWidget(self)
                self.drag_drop_layout.addWidget(self.drag_drop_widget, 4)
                self.drag_drop_layout.addWidget(self.file_selection_display_label, 1,
                                                alignment=Qt.AlignmentFlag.AlignCenter)
                self.file_selection_group_box.setLayout(self.drag_drop_layout)

                # File Tree
                self.file_tree = QTreeWidget()
                self.file_tree_layout = QHBoxLayout()
                self.file_tree.setContextMenuPolicy(Qt.ContextMenuPolicy.CustomContextMenu)
                self.file_tree.customContextMenuRequested.connect(self.open_context_menu)
                self.file_tree.setHeaderLabels(["Name", "Type", "Size"])
                self.file_tree_layout.addWidget(self.file_tree)
                self.file_tree.setStyleSheet(self.QTree_stylesheet)
                self.file_view_group_box.setLayout(self.file_tree_layout)

                self.fileUpload_layout.addWidget(self.file_selection_group_box, 1)
                self.fileUpload_layout.addWidget(self.file_view_group_box, 1)
                self.fileupload_container = QWidget(self)
                self.fileupload_container.setLayout(self.fileUpload_layout)
                self.fileupload_container.setFixedSize(1150, 300)

                # Figure Layout
                self.figure_layout = QHBoxLayout()

                # Raw Canvas
                self.raw_canvas_layout = QVBoxLayout()
                self.raw_canvas = MplCanvas(self, width=6, height=5, dpi=100)
                self.raw_canvas.ax2 = self.raw_canvas.ax.twinx()
                self.raw_canvas.ax2.tick_params(right=False, labelright=False)
                self.raw_toolbar = NavigationToolbar(self.raw_canvas, self)
                self.raw_toolbar.setStyleSheet("QWidget { border: None; }")
                self.raw_canvas.ax.set_title("Hysteresis Loop")
                self.raw_canvas_layout.addWidget(self.raw_toolbar, alignment=Qt.AlignmentFlag.AlignCenter)
                self.raw_canvas_layout.addWidget(self.raw_canvas)

                # Fit Canvas
                self.fit_canvas_layout = QVBoxLayout()
                self.fit_canvas = MplCanvas(self, width=6, height=5, dpi=100)
                self.fit_toolbar = NavigationToolbar(self.fit_canvas, self)
                self.fit_toolbar.setStyleSheet("QWidget { border: None; }")
                self.fit_canvas.ax.set_title("Fitting")
                self.fit_canvas_layout.addWidget(self.fit_toolbar, alignment=Qt.AlignmentFlag.AlignCenter)
                self.fit_canvas_layout.addWidget(self.fit_canvas)

                # Summary Canvas
                self.summary_canvas_layout = QVBoxLayout()
                self.summary_canvas = MplCanvas(self, width=6, height=5, dpi=100)
                self.summary_toolbar = NavigationToolbar(self.summary_canvas, self)
                self.summary_toolbar.setStyleSheet("QWidget { border: None; }")
                self.summary_canvas.ax.set_title("Fitting Summary")
                self.summary_canvas_layout.addWidget(self.summary_toolbar, alignment=Qt.AlignmentFlag.AlignCenter)
                self.summary_canvas_layout.addWidget(self.summary_canvas)

                self.figure_layout.addLayout(self.raw_canvas_layout)
                self.figure_layout.addLayout(self.fit_canvas_layout)
                self.figure_layout.addLayout(self.summary_canvas_layout)
                self.figure_container = QWidget(self)
                self.figure_container.setLayout(self.figure_layout)
                self.figure_container.setFixedSize(1150, 400)

                # Fitting Parameters
                self.fitting_parameter_layout = QHBoxLayout()
                self.plot_control_group_box = QGroupBox("Plot Control")
                self.fit_control_group_box = QGroupBox("Fit Control")
                self.fit_summary_group_box = QGroupBox("Fit Summary")

                self.plot_control_layout = QVBoxLayout()
                self.fit_control_layout = QVBoxLayout()
                self.fit_summary_layout = QVBoxLayout()

                # Plot Control Checkboxes
                self.raw_plot_check_box = QCheckBox("Show Raw Data")
                self.processed_plot_check_box = QCheckBox("Show Processed Data")
                self.area_plot_check_box = QCheckBox("Show Area")
                self.x_correction_check_box = QCheckBox("Show Hc corrected plot")
                self.y_correction_check_box = QCheckBox("Show Ms corrected plot")
                self.xy_correction_check_box = QCheckBox("Show both corrected plot")
                self.thermal_drift_check_box = QCheckBox("Show Thermal Drift plot")
                self.slope_check_box = QCheckBox("Show slope")

                self.plot_control_layout.addWidget(self.raw_plot_check_box)
                self.plot_control_layout.addWidget(self.slope_check_box)
                self.plot_control_layout.addWidget(self.area_plot_check_box)
                self.plot_control_layout.addWidget(self.processed_plot_check_box)
                self.plot_control_layout.addWidget(self.x_correction_check_box)
                self.plot_control_layout.addWidget(self.y_correction_check_box)
                self.plot_control_layout.addWidget(self.xy_correction_check_box)
                self.plot_control_layout.addWidget(self.thermal_drift_check_box)
                self.plot_control_group_box.setLayout(self.plot_control_layout)

                # Fit Control Entries
                self.area_fit_results_layout = QHBoxLayout()
                self.area_text_label = QLabel('Area: ')
                self.area_entry = QLineEdit()
                self.area_entry.setReadOnly(True)
                self.area_fit_results_layout.addWidget(self.area_text_label)
                self.area_fit_results_layout.addWidget(self.area_entry)

                self.Hc_fit_results_layout = QHBoxLayout()
                self.hc_text_label = QLabel('Coercivity: ')
                self.hc_entry = QLineEdit()
                self.hc_entry.setReadOnly(True)
                self.Hc_fit_results_layout.addWidget(self.hc_text_label)
                self.Hc_fit_results_layout.addWidget(self.hc_entry)

                self.ms_fit_results_layout = QHBoxLayout()
                self.ms_text_label = QLabel('Saturation Magnetization: ')
                self.ms_entry = QLineEdit()
                self.ms_entry.setReadOnly(True)
                self.ms_fit_results_layout.addWidget(self.ms_text_label)
                self.ms_fit_results_layout.addWidget(self.ms_entry)

                self.eb_fit_results_layout = QHBoxLayout()
                self.eb_text_label = QLabel('Exchange Bias: ')
                self.eb_entry = QLineEdit()
                self.eb_entry.setReadOnly(True)
                self.eb_fit_results_layout.addWidget(self.eb_text_label)
                self.eb_fit_results_layout.addWidget(self.eb_entry)

                self.fit_control_layout.addLayout(self.area_fit_results_layout)
                self.fit_control_layout.addLayout(self.Hc_fit_results_layout)
                self.fit_control_layout.addLayout(self.ms_fit_results_layout)
                self.fit_control_layout.addLayout(self.eb_fit_results_layout)
                self.fit_control_group_box.setLayout(self.fit_control_layout)

                # Summary Checkboxes
                self.area_check_box = QCheckBox("Area")
                self.coercivity_check_box = QCheckBox("Coercivity")
                self.ms_check_box = QCheckBox("Saturation Magnetization")
                self.eb_check_box = QCheckBox("Exchange Bias")

                self.fit_summary_layout.addWidget(self.area_check_box)
                self.fit_summary_layout.addWidget(self.coercivity_check_box)
                self.fit_summary_layout.addWidget(self.ms_check_box)
                self.fit_summary_layout.addWidget(self.eb_check_box)
                self.fit_summary_group_box.setLayout(self.fit_summary_layout)

                self.fitting_parameter_layout.addWidget(self.plot_control_group_box)
                self.fitting_parameter_layout.addWidget(self.fit_control_group_box)
                self.fitting_parameter_layout.addWidget(self.fit_summary_group_box)

                self.plot_control_group_box.setFixedSize(380, 200)
                self.fit_control_group_box.setFixedSize(380, 200)
                self.fit_summary_group_box.setFixedSize(380, 200)

                # Add Process and Reset Buttons
                self.button_layout = QHBoxLayout()
                self.process_button = QPushButton("Process Data")
                self.process_button.setStyleSheet(self.Button_stylesheet)
                self.process_button.clicked.connect(self.start_processing)
                self.process_button.setFixedSize(200, 40)

                self.reset_button = QPushButton("Reset")
                self.reset_button.setStyleSheet(self.Button_stylesheet)
                self.reset_button.clicked.connect(self.reset_all)
                self.reset_button.setFixedSize(200, 40)

                self.button_layout.addStretch()
                self.button_layout.addWidget(self.process_button)
                self.button_layout.addWidget(self.reset_button)

                self.button_container = QWidget(self)
                self.button_container.setLayout(self.button_layout)
                self.button_container.setFixedSize(1150, 50)

                # Add Progress Bar (initially hidden)
                self.progress_layout = QHBoxLayout()
                self.progress_label = QLabel("Processing data...")
                self.progress_label.setStyleSheet("font-size: 12px; color: #4b6172;")
                self.progress_bar = QProgressBar()
                self.progress_bar.setStyleSheet("""
                    QProgressBar {
                        border: 2px solid grey;
                        border-radius: 5px;
                        text-align: center;
                        height: 25px;
                    }
                    QProgressBar::chunk {
                        background-color: #4CAF50;
                        width: 10px;
                    }
                """)
                self.progress_bar.setFixedWidth(750)
                self.progress_layout.addWidget(self.progress_label)
                self.progress_layout.addWidget(self.progress_bar)
                self.progress_layout.addStretch()

                # Hide progress bar initially
                self.progress_label.hide()
                self.progress_bar.hide()

                # Connect signals
                self.raw_plot_check_box.stateChanged.connect(self.update_plots)
                self.slope_check_box.stateChanged.connect(self.update_plots)
                self.processed_plot_check_box.stateChanged.connect(self.update_plots)
                self.area_plot_check_box.stateChanged.connect(self.update_plots)
                self.x_correction_check_box.stateChanged.connect(self.update_plots)
                self.y_correction_check_box.stateChanged.connect(self.update_plots)
                self.xy_correction_check_box.stateChanged.connect(self.update_plots)
                self.thermal_drift_check_box.stateChanged.connect(self.update_plots)

                self.area_check_box.stateChanged.connect(self.update_summary_plot)
                self.coercivity_check_box.stateChanged.connect(self.update_summary_plot)
                self.ms_check_box.stateChanged.connect(self.update_summary_plot)
                self.eb_check_box.stateChanged.connect(self.update_summary_plot)

                self.file_tree.itemSelectionChanged.connect(self.on_item_selection_changed)

                # Add to main layout
                self.VSM_data_extraction_main_layout.addWidget(self.VSM_QD_EXTRACT_label,
                                                               alignment=Qt.AlignmentFlag.AlignTop)
                self.VSM_data_extraction_main_layout.addWidget(self.fileupload_container)
                self.VSM_data_extraction_main_layout.addWidget(self.button_container)
                self.VSM_data_extraction_main_layout.addLayout(self.progress_layout)
                self.VSM_data_extraction_main_layout.addWidget(self.figure_container)
                self.VSM_data_extraction_main_layout.addLayout(self.fitting_parameter_layout)

                self.VSM_data_extraction_main_layout.addStretch(1)

                self.setCentralWidget(self.scroll_area)

                print("UI initialization complete!")

        except Exception as e:
            print(f"Error in init_ui: {e}")
            import traceback
            traceback.print_exc()
            QMessageBox.warning(self, "Error", f"Failed to initialize UI: {str(e)}")
            return

    def display_files(self, folder_path):
        # Check if UI is initialized
        if not hasattr(self, 'file_tree') or not hasattr(self, 'file_selection_display_label'):
            QMessageBox.warning(self, "Error", "UI not properly initialized. Please restart the application.")
            return

        self.file_tree.clear()
        self.file_in_list.clear()
        self.folder = folder_path
        self.folder_selected = folder_path

        # Reset processed flag when new files are loaded
        self.is_processed = False

        # Initialize folder structure
        self.ProcessedRAW = self.folder_selected + 'Processed_Graph'
        self.ProcCal = self.folder_selected + 'Processed_Result'

        self.file_selection_display_label.setText("Directory Successfully Uploaded")
        self.file_selection_display_label.setStyleSheet("""
            color: #4b6172; font-size: 12px; background-color: #DfE7Ef; 
            border-radius: 5px; padding: 5px;
        """)

        # Check if Processed_Data subfolder exists, if not, look in root
        processed_data_path = os.path.join(folder_path, 'Processed_Data')
        if os.path.exists(processed_data_path):
            search_path = processed_data_path
        else:
            # Look for CSV files directly in the uploaded folder
            search_path = folder_path

        for root, dirs, files in os.walk(search_path):
            for file_name in files:
                if file_name.endswith('.csv'):
                    file_path = os.path.join(root, file_name)
                    self.file_in_list.append(file_path)
                    file_info = os.stat(file_path)
                    file_size_kb = file_info.st_size / 1024
                    file_size_str = f"{file_size_kb:.2f} KB"
                    file_type = 'csv'
                    item = QTreeWidgetItem(self.file_tree, [file_name, file_type, file_size_str, ""])
                    item.setToolTip(0, file_path)

        if len(self.file_in_list) == 0:
            QMessageBox.warning(self, "No CSV Files",
                                "No CSV files found. Please upload a folder containing CSV files with VSM data.\n\n" +
                                "Expected format:\n" +
                                "- Files named like: 10K.csv, 50K.csv, 300K.csv\n" +
                                "- Each CSV should have columns: Temperature (K), Magnetic Field (Oe), Moment (emu)")
            return

        # Don't auto-process, let user click the button
        self.file_tree.resizeColumnToContents(0)
        self.file_tree.resizeColumnToContents(1)
        self.file_tree.resizeColumnToContents(2)

    def display_multiple_files(self, file_paths):
        # Initialize folder from first file if not set
        if not self.folder_selected and file_paths:
            first_file_dir = os.path.dirname(file_paths[0])
            self.folder_selected = first_file_dir + '/'
            self.ProcessedRAW = self.folder_selected + 'Processed_Graph'
            self.ProcCal = self.folder_selected + 'Processed_Result'

        # Reset processed flag
        self.is_processed = False

        current_files = {self.file_tree.topLevelItem(i).text(0): self.file_tree.topLevelItem(i)
                         for i in range(self.file_tree.topLevelItemCount())}

        for file_path in file_paths:
            file_name = os.path.basename(file_path)
            if file_name not in current_files:
                self.file_in_list.append(file_path)
                file_info = os.stat(file_path)
                file_size_kb = file_info.st_size / 1024
                file_size_str = f"{file_size_kb:.2f} KB"
                file_type = 'csv'
                item = QTreeWidgetItem(self.file_tree, [file_name, file_type, file_size_str, ""])
                item.setToolTip(0, file_path)

        # Don't auto-process
        self.file_tree.resizeColumnToContents(0)
        self.file_tree.resizeColumnToContents(1)
        self.file_tree.resizeColumnToContents(2)

        self.file_selection_display_label.setText(f"{self.file_tree.topLevelItemCount()} Files Successfully Uploaded")
        self.file_selection_display_label.setStyleSheet("""
            color: #4b6172; font-size: 12px; background-color: #DfE7Ef; 
            border-radius: 5px; padding: 5px;
        """)

    def on_item_selection_changed(self):
        """Handle file selection and update all plots"""
        if not self.is_processed:
            # Clear selection if not processed
            self.file_tree.clearSelection()
            return

        selected_items = self.file_tree.selectedItems()
        if selected_items:
            selected_item = selected_items[0]
            file_name = selected_item.text(0)  # Get filename from tree
            file_path = selected_item.toolTip(0)  # Get full path from tooltip

            # Load and display the selected file
            self.load_and_display_file_from_tree(file_name, file_path)

    def load_and_display_file_from_tree(self, file_name, file_path):
        """Load file data from tree selection and display in plots"""
        try:
            # Extract temperature from filename
            cur_temp = None
            # for i in range(len(file_name) - 1, 0, -1):
            #     if file_name[i] == 'K':
            #         cur_temp = file_name[:i]
            #         break

            hysteresis_temp_df = pd.read_csv(file_path, header=0, engine='c')
            thermal = hysteresis_temp_df.iloc[:, 0]
            cur_temp = round(thermal.mean(), 1)
            cur_temp_file = str(cur_temp).replace(".", "_")
            # cur_temp_file = cur_temp_file[:cur_temp_file.rfind('_')] if cur_temp_file[-2:] == '_0' else cur_temp_file
            if cur_temp_file.endswith('_0'):
                cur_temp_file = cur_temp_file[:cur_temp_file.rfind('_')]

            # Case 2: ends with "_9" → parse the part before "_" as int, then +1
            elif cur_temp_file.endswith('_9'):
                head = cur_temp_file[:cur_temp_file.rfind('_')]
                try:
                    cur_temp_file = str(int(head) + 1)
                except ValueError:
                    # Fallback if the part before '_' isn't a clean integer
                    # (optional) leave unchanged or handle/log as needed
                    pass

            if cur_temp is None:
                QMessageBox.warning(self, "Error",
                                    f"Cannot extract temperature from filename: {file_name}\nFiles must be named like: 10K.csv, 300K.csv")
                return

            # Initialize data storage
            self.current_file_data = {
                'temperature': cur_temp,
                'file_name': file_name,
                'raw_data': None,
                'processed_data': None,
                'processed_raw_data': None,
                'split_data': None,
                'thermal_drift': None,
                'slope': None,
                'raw_data_fit': None,
                'ms_corrected': None,
                'eb_corrected': None,
                'eb_ms_corrected': None,
                'split_data_fit': None
            }

            # Load data directly from the file in tree view
            if os.path.exists(file_path):
                try:
                    raw_df = pd.read_csv(file_path, header=0)
                    if len(raw_df.columns) >= 3:
                        self.current_file_data['raw_data'] = {
                            'x': raw_df.iloc[:, 1].values,
                            'y': raw_df.iloc[:, 2].values
                        }

                    else:
                        QMessageBox.warning(self, "Error",
                                            f"File must have at least 3 columns.\nFound: {len(raw_df.columns)} columns")
                        return
                except Exception as e:
                    QMessageBox.warning(self, "Error Loading File", f"Cannot read file: {file_name}\n\nError: {str(e)}")
                    return
            else:
                QMessageBox.warning(self, "File Not Found", f"File not found: {file_path}")
                return

            # Check for processed data only if ProcCal exists
            if self.ProcCal and os.path.exists(self.ProcCal):
                processed_data_path = os.path.join(self.ProcCal, 'Slope_Removal', f'{cur_temp_file}K_slope_removal_data.csv')
                processed_raw_path = os.path.join(self.ProcCal, 'Final_Processed_RAW_Data',
                                                  f'{cur_temp_file}K_Final_RAW.csv')
                split_data_path = os.path.join(self.ProcCal,'Slope_Removal_Spilt', f'{cur_temp_file}K_slope_removal_spilt_data.csv')
                slope_data_path = os.path.join(self.ProcCal, 'Slope', f'{cur_temp_file}K_slope_data.csv')
                thermal_data_path = os.path.join(self.ProcCal, 'Thermal_Difference', f'{cur_temp_file}K_thermal_difference.csv')
                raw_data_fit_path = os.path.join(self.ProcCal, f'{cur_temp_file}K_raw_fit.csv')
                ms_corrected_data_path = os.path.join(self.ProcCal, 'Ms_Corrected',f'{cur_temp_file}K_ms_corrected_data.csv')
                eb_corrected_data_path = os.path.join(self.ProcCal, 'Eb_Corrected', f'{cur_temp_file}K_eb_corrected_data.csv')
                eb_ms_corrected_data_path = os.path.join(self.ProcCal, 'Ms_Eb_Corrected', f'{cur_temp_file}K_eb_ms_corrected_data.csv')
                fitted_curve_data_path = os.path.join(self.ProcCal, 'Fitted_Curve',
                                                         f'{cur_temp_file}K_ms_corrected_data_fitting_curve_upper_curve.csv')
                # Load processed thermal drift data if available
                if os.path.exists(thermal_data_path):
                    try:
                        thermal_drift_df = pd.read_csv(thermal_data_path, header=None)
                        self.current_file_data['thermal_drift'] = {
                            'x': thermal_drift_df.iloc[:, 0].values,
                            'y': thermal_drift_df.iloc[:, 1].values
                        }
                    except:
                        pass

                if os.path.exists(fitted_curve_data_path):
                    try:
                        raw_fit_df = pd.read_csv(raw_data_fit_path, header=None)
                        self.current_file_data['raw_data_fit'] = {
                            'x1': raw_fit_df.iloc[:, 0].values,
                            'y1': raw_fit_df.iloc[:, 1].values,
                            'x2': raw_fit_df.iloc[:, 2].values,
                            'y2': raw_fit_df.iloc[:, 3].values
                        }
                    except:
                        pass

                # Load processed data if available
                if os.path.exists(processed_data_path):
                    try:
                        proc_df = pd.read_csv(processed_data_path, header=None)
                        self.current_file_data['processed_data'] = {
                            'x': proc_df.iloc[:, 0].values,
                            'y': proc_df.iloc[:, 1].values
                        }
                    except:
                        pass

                # Load eb corrected data if available
                if os.path.exists(eb_corrected_data_path):
                    try:
                        proc_df = pd.read_csv(eb_corrected_data_path, header=None)
                        self.current_file_data['eb_corrected'] = {
                            'x': proc_df.iloc[:, 0].values,
                            'y': proc_df.iloc[:, 1].values
                        }
                    except:
                        pass

                # Load ms corrected data if available
                if os.path.exists(ms_corrected_data_path):
                    try:
                        proc_df = pd.read_csv(ms_corrected_data_path, header=None)
                        self.current_file_data['ms_corrected'] = {
                            'x': proc_df.iloc[:, 0].values,
                            'y': proc_df.iloc[:, 1].values
                        }
                    except:
                        pass

                # Load ms eb corrected data if available
                if os.path.exists(eb_ms_corrected_data_path):
                    try:
                        proc_df = pd.read_csv(eb_ms_corrected_data_path, header=None)
                        self.current_file_data['eb_ms_corrected'] = {
                            'x': proc_df.iloc[:, 0].values,
                            'y': proc_df.iloc[:, 1].values
                        }
                    except:
                        pass

                # Load processed raw data if available
                if os.path.exists(processed_raw_path):
                    try:
                        proc_raw_df = pd.read_csv(processed_raw_path, header=None)
                        self.current_file_data['processed_raw_data'] = {
                            'x': proc_raw_df.iloc[:, 0].values,
                            'y': proc_raw_df.iloc[:, 1].values
                        }
                    except:
                        pass

                # Load split data if available
                if os.path.exists(split_data_path):
                    try:
                        split_df = pd.read_csv(split_data_path, header=None)
                        self.current_file_data['split_data'] = {
                            'x_lower': split_df.iloc[:, 0].dropna().values,
                            'y_lower': split_df.iloc[:, 1].dropna().values,
                            'x_upper': split_df.iloc[:, 2].dropna().values,
                            'y_upper': split_df.iloc[:, 3].dropna().values
                        }
                    except:
                        pass

                # Load slope data if available
                if os.path.exists(slope_data_path):
                    try:
                        slope_df = pd.read_csv(slope_data_path, header=None)
                        self.current_file_data['slope'] = {
                            'x': slope_df.iloc[:, 0].values,
                            'y': slope_df.iloc[:, 1].values
                        }
                    except:
                        pass
            # Update displays with whatever data is available
            self.update_plots()
            self.update_fitting_display()

        except Exception as e:
            QMessageBox.warning(self, "Error", f"Unexpected error: {str(e)}")

    def load_all_files_from_tree(self):
        """Load all files from tree view for batch processing"""
        all_files = []
        root = self.file_tree.invisibleRootItem()
        child_count = root.childCount()

        for i in range(child_count):
            item = root.child(i)
            file_name = item.text(0)
            file_path = item.toolTip(0)
            all_files.append({
                'name': file_name,
                'path': file_path,
                'type': item.text(1),
                'size': item.text(2)
            })

        return all_files

    def update_plots(self):
        """Update the hysteresis and fitting plots based on checkbox states"""
        if self.current_file_data is None:
            return

        # Clear canvases
        self.raw_canvas.ax.clear()
        self.raw_canvas.ax2.clear()
        self.raw_canvas.ax2.tick_params(right=False, labelright=False)
        self.fit_canvas.ax.clear()

        # Get checkbox states
        show_raw = self.raw_plot_check_box.isChecked()
        show_slope = self.slope_check_box.isChecked()
        show_processed = self.processed_plot_check_box.isChecked()
        show_area = self.area_plot_check_box.isChecked()
        show_x_corr = self.x_correction_check_box.isChecked()
        show_y_corr = self.y_correction_check_box.isChecked()
        show_xy_corr = self.xy_correction_check_box.isChecked()
        show_thermal_drift = self.thermal_drift_check_box.isChecked()

        # Plot on raw canvas
        if show_raw and self.current_file_data['raw_data'] is not None:
            self.raw_canvas.ax.scatter(
                self.current_file_data['raw_data']['x'],
                self.current_file_data['raw_data']['y'],
                s=1, alpha=0.6, color='tomato', label='Raw Data'
            )

        if show_y_corr and self.current_file_data['ms_corrected'] is not None:
            self.raw_canvas.ax.scatter(
                self.current_file_data['ms_corrected']['x'],
                self.current_file_data['ms_corrected']['y'],
                s=1, alpha=0.6, color='blue', label='Raw Fit Data'
            )

        if show_slope and self.current_file_data['slope'] is not None:
            self.raw_canvas.ax.scatter(
                self.current_file_data['slope']['x'],
                self.current_file_data['slope']['y'],
                s=1, alpha=0.6, color='sienna', label='Slope Data'
            )

        if show_processed and self.current_file_data['processed_data'] is not None:
            self.raw_canvas.ax.scatter(
                self.current_file_data['processed_data']['x'],
                self.current_file_data['processed_data']['y'],
                s=1, alpha=0.6, color='black', label='Processed Data'
            )

        if show_x_corr and self.current_file_data['eb_corrected'] is not None:
            self.raw_canvas.ax.scatter(
                self.current_file_data['eb_corrected']['x'],
                self.current_file_data['eb_corrected']['y'],
                s=1, alpha=0.6, color='green', label='Hc Corrected'
            )

        if show_xy_corr and self.current_file_data['eb_ms_corrected'] is not None:
            self.raw_canvas.ax.scatter(
                self.current_file_data['eb_ms_corrected']['x'],
                self.current_file_data['eb_ms_corrected']['y'],
                s=1, alpha=0.6, color='red', label='Both Corrected'
            )

        if show_area and self.current_file_data['raw_data'] is not None:
            x = self.current_file_data['raw_data']['x']
            y = self.current_file_data['raw_data']['y']
            df_area = pd.concat([pd.Series(x, name='Field'), pd.Series(y, name='Moment')], axis=1)
            df_area_drop_nan = df_area.dropna()
            x_drop = df_area_drop_nan.iloc[:, 0]
            y_drop = df_area_drop_nan.iloc[:, 1]
            self.raw_canvas.ax.fill_between(x_drop, y_drop, alpha=0.3, label='Area')

        if show_thermal_drift and self.current_file_data['thermal_drift'] is not None:
            self.raw_canvas.ax2.scatter(
                self.current_file_data['thermal_drift']['x'],
                self.current_file_data['thermal_drift']['y'],
                s=1, alpha=0.4, color='grey', label='Raw Data'
            )
            self.raw_canvas.ax2.set_ylabel('T (K)', fontsize=11)
            self.raw_canvas.ax2.yaxis.set_label_position("right")
            self.raw_canvas.ax2.tick_params(right=True, labelright=True,labelsize=9)
            self.raw_canvas.ax2.ticklabel_format(style='plain', axis='y')
            self.raw_canvas.ax2.yaxis.get_offset_text().set_x(1.05)

        self.raw_canvas.ax.set_xlabel('Magnetic Field (Oe)', fontsize=11)
        self.raw_canvas.ax.set_ylabel('Moment (emu)', fontsize=11)
        self.raw_canvas.ax.set_title(f"{self.current_file_data['temperature']} K Hysteresis Loop", fontsize=12)
        if show_raw or show_processed or show_area or show_x_corr or show_y_corr or show_xy_corr:
            handles, labels = self.raw_canvas.ax.get_legend_handles_labels()
            if handles:
                self.raw_canvas.ax.legend(fontsize=8, loc='best')
        self.raw_canvas.ax.grid(True, alpha=0.3)
        self.raw_canvas.ax.tick_params(labelsize=9)
        self.raw_canvas.fig.tight_layout()
        self.raw_canvas.draw()

        # Plot on fit canvas
        if self.current_file_data['split_data'] is not None:
            split = self.current_file_data['split_data']

            self.fit_canvas.ax.scatter(split['x_lower'], split['y_lower'],
                                       s=10, alpha=0.2, color='cornflowerblue', label='Lower Branch')
            self.fit_canvas.ax.scatter(split['x_upper'], split['y_upper'],
                                       s=10, alpha=0.2, color='bisque', label='Upper Branch')

            # Add fitted curves
            try:
                upper_fit = extract_coercivity_method1_fit(split['x_lower'], split['y_lower'], include_slope=False)
                lower_fit = extract_coercivity_method1_fit(split['x_upper'], split['y_upper'], include_slope=False)
                if upper_fit['success'] and lower_fit['success']:
                    self.fit_canvas.ax.plot(split['x_lower'], upper_fit['M_fit'], color='mediumblue', linestyle='-',
                                            linewidth=2, label='Lower Fit')
                    self.fit_canvas.ax.plot(split['x_upper'], lower_fit['M_fit'], color='darkorange', linestyle='-',
                                            linewidth=2, label='Upper Fit')


            except Exception as e:
                print(f"Fitting error: {e}")

        self.fit_canvas.ax.set_xlabel('Magnetic Field (Oe)', fontsize=11)
        self.fit_canvas.ax.set_ylabel('Moment (emu)', fontsize=11)
        self.fit_canvas.ax.set_title('Hysteresis Fitting', fontsize=12)
        handles, labels = self.fit_canvas.ax.get_legend_handles_labels()
        if handles:
            self.fit_canvas.ax.legend(fontsize=8, loc='best')
        self.fit_canvas.ax.grid(True, alpha=0.3)
        self.fit_canvas.ax.tick_params(labelsize=9)
        self.fit_canvas.fig.tight_layout()
        self.fit_canvas.draw()

    def update_fitting_display(self):
        """Update the fitting parameter displays"""
        if self.current_file_data is None:
            return

        temp = self.current_file_data['temperature']

        # Load and display area
        if self.area_df is not None:
            # temp_match = self.area_df[self.area_df['Temperature'].astype(str) == temp]
            temp_match = self.area_df[self.area_df['Temperature'] == temp]
            if not temp_match.empty:
                area_val = temp_match['Area'].values[0]
                self.area_entry.setText(f"{area_val:.6e}")
            else:
                self.area_entry.setText("N/A")

        # Load and display coercivity
        if self.Coercivity_df is not None:
            # temp_match = self.Coercivity_df[self.Coercivity_df['Temperature'].astype(str) == temp]
            temp_match = self.Coercivity_df[self.Coercivity_df['Temperature'] == temp]
            if not temp_match.empty:
                hc_val = temp_match['Coercivity'].values[0]
                self.hc_entry.setText(f"{hc_val:.4f}")
            else:
                self.hc_entry.setText("N/A")

        # Load and display Ms
        if self.Ms_df is not None:
            # temp_match = self.Ms_df[self.Ms_df['Temperature'].astype(str) == temp]
            temp_match = self.Ms_df[self.Ms_df['Temperature'] == temp]
            if not temp_match.empty:
                ms_val = temp_match['Saturation Field'].values[0]
                self.ms_entry.setText(f"{ms_val:.6e}")
            else:
                self.ms_entry.setText("N/A")

        # Exchange bias
        if self.eb_df is not None:
            # temp_match = self.eb_df[self.eb_df['Temperature'].astype(str) == temp]
            temp_match = self.eb_df[self.eb_df['Temperature'] == temp]
            if not temp_match.empty:
                eb_val = temp_match['Exchange Bias'].values[0]
                self.eb_entry.setText(f"{eb_val:.4f}")
            else:
                self.eb_entry.setText("N/A")


    def update_summary_plot(self):
        """Update the summary plot based on selected parameters"""
        self.summary_canvas.fig.clear()

        # Check which parameters to plot
        plot_area = self.area_check_box.isChecked()
        plot_coercivity = self.coercivity_check_box.isChecked()
        plot_ms = self.ms_check_box.isChecked()
        plot_eb = self.eb_check_box.isChecked()

        num_plots = sum([plot_area, plot_coercivity, plot_ms, plot_eb])

        if num_plots == 0:
            ax = self.summary_canvas.fig.add_subplot(111)
            ax.text(0.5, 0.5, 'Select parameters to plot',
                    ha='center', va='center', fontsize=14)
            ax.set_title("Fitting Summary", fontsize=12)

            self.summary_canvas.draw()
            return

        # Plot selected parameters
        plot_index = 1

        if plot_area and self.area_df is not None:
            if num_plots > 1:
                ax = self.summary_canvas.fig.add_subplot(2, 2, plot_index)
            else:
                ax = self.summary_canvas.fig.add_subplot(111)

            ax.plot(self.area_df['Temperature'], self.area_df['Area'],
                    'o-', color='blue', markersize=5, linewidth=2)
            ax.set_xlabel('Temperature (K)', fontsize=11)
            ax.set_ylabel('Area', fontsize=11)
            ax.set_title("Area", fontsize=12)
            ax.grid(True, alpha=0.3)
            plot_index += 1

        if plot_coercivity and self.Coercivity_df is not None:
            if num_plots > 1:
                ax = self.summary_canvas.fig.add_subplot(2, 2, plot_index)
            else:
                ax = self.summary_canvas.fig.add_subplot(111)

            ax.plot(self.Coercivity_df['Temperature'], self.Coercivity_df['Coercivity'],
                    'o-', color='red', markersize=5, linewidth=2)
            ax.set_xlabel('Temperature (K)', fontsize=11)
            ax.set_ylabel('Coercivity (Oe)', fontsize=11)
            ax.set_title("Coercivity", fontsize=12)
            ax.grid(True, alpha=0.3)
            plot_index += 1

        if plot_ms and self.Ms_df is not None:
            if num_plots > 1:
                ax = self.summary_canvas.fig.add_subplot(2, 2, plot_index)
            else:
                ax = self.summary_canvas.fig.add_subplot(111)

            ax.plot(self.Ms_df['Temperature'], self.Ms_df['Saturation Field'],
                    'o-', color='green', markersize=5, linewidth=2)
            ax.set_xlabel('Temperature (K)', fontsize=11)
            ax.set_ylabel('Ms (emu)', fontsize=11)
            ax.set_title("Ms", fontsize=12)
            ax.grid(True, alpha=0.3)
            plot_index += 1

        if plot_eb and self.eb_df is not None:
            if num_plots > 1:
                ax = self.summary_canvas.fig.add_subplot(2, 2, plot_index)
            else:
                ax = self.summary_canvas.fig.add_subplot(111)

            ax.plot(self.eb_df['Temperature'], self.eb_df['Exchange Bias'],
                    'o-', color='purple', markersize=5, linewidth=2)
            ax.set_xlabel('Temperature (K)', fontsize=10)
            ax.set_ylabel('Exchange Bias (Oe)', fontsize=10)
            ax.grid(True, alpha=0.3)

        self.summary_canvas.fig.tight_layout()
        self.summary_canvas.draw()

    def start_processing(self):
        """Start the data processing with progress bar"""
        if self.file_tree.topLevelItemCount() == 0:
            QMessageBox.warning(self, "No Files", "Please upload CSV files first.")
            return

        # Show progress bar
        self.progress_label.show()
        self.progress_bar.show()
        self.progress_bar.setValue(0)
        self.process_button.setEnabled(False)

        # Process data
        try:
            self.save_folder_path = QFileDialog.getExistingDirectory(self, "Select Folder")
            print(self.save_folder_path)
            self.process_data()
        except Exception as e:
            tb_str = traceback.format_exc()
            QMessageBox.warning(self, "Processing Error", f"Error during processing:\n{str(e)} \n{str(tb_str)}")
        finally:
            # Hide progress bar
            self.progress_label.hide()
            self.progress_bar.hide()
            self.process_button.setEnabled(True)

    def reset_all(self):
        """Reset all data and clear the interface"""
        reply = QMessageBox.question(self, 'Reset Confirmation',
                                     'Are you sure you want to reset? This will clear all uploaded files and processed data.',
                                     QMessageBox.StandardButton.Yes | QMessageBox.StandardButton.No,
                                     QMessageBox.StandardButton.No)

        if reply == QMessageBox.StandardButton.Yes:
            # Clear tree
            self.file_tree.clear()
            self.file_in_list.clear()

            # Clear plots
            self.raw_canvas.ax.clear()
            self.raw_canvas.ax2.clear()
            self.raw_canvas.ax2.tick_params(right=False, labelright=False)
            self.raw_canvas.ax.set_title("Hysteresis Loop")
            self.raw_canvas.draw()

            self.fit_canvas.ax.clear()
            self.fit_canvas.ax.set_title("Fitting")
            self.fit_canvas.draw()

            self.summary_canvas.ax.clear()
            self.summary_canvas.ax.set_title("Fitting Summary")
            self.summary_canvas.draw()

            # Clear parameter fields
            self.area_entry.clear()
            self.hc_entry.clear()
            self.ms_entry.clear()
            self.eb_entry.clear()

            # Clear checkboxes
            self.raw_plot_check_box.setChecked(False)
            self.processed_plot_check_box.setChecked(False)
            self.area_plot_check_box.setChecked(False)
            self.x_correction_check_box.setChecked(False)
            self.y_correction_check_box.setChecked(False)
            self.xy_correction_check_box.setChecked(False)
            self.thermal_drift_check_box.setChecked(False)
            self.area_check_box.setChecked(False)
            self.coercivity_check_box.setChecked(False)
            self.ms_check_box.setChecked(False)
            self.eb_check_box.setChecked(False)

            # Reset data
            self.current_file_data = None
            self.area_df = None
            self.Ms_df = None
            self.Coercivity_df = None
            self.eb_df = None
            self.folder_selected = None
            self.save_folder_path = None
            self.ProcCal = None
            self.ProcessedRAW = None
            self.is_processed = False

            # Disable selection
            self.file_tree.setSelectionMode(QTreeWidget.SelectionMode.NoSelection)

            # Reset label
            self.file_selection_display_label.setText('Please Upload Files or Directory')
            self.file_selection_display_label.setStyleSheet("""
                color: white; font-size: 12px; background-color: #f38d76; 
                border-radius: 5px; padding: 5px;
            """)

    def open_context_menu(self, position: QPoint):
        """Open the context menu on right-click"""
        menu = QMenu()
        remove_action = QWidgetAction(self)
        remove_label = QLabel("Remove")
        remove_label.mousePressEvent = lambda event: self.handle_remove_click(event)
        remove_action.setDefaultWidget(remove_label)
        menu.addAction(remove_action)
        menu.exec(self.file_tree.viewport().mapToGlobal(position))

    def handle_remove_click(self, event):
        """Handle right-click on remove label"""
        if event.button() == Qt.MouseButton.LeftButton:
            self.remove_selected_item()

    def remove_selected_item(self):
        """Remove the selected item from the tree"""
        selected_item = self.file_tree.currentItem()
        if selected_item:
            file_path = selected_item.toolTip(0)
            if file_path in self.file_in_list:
                self.file_in_list.remove(file_path)
            index = self.file_tree.indexOfTopLevelItem(selected_item)
            if index != -1:
                self.file_tree.takeTopLevelItem(index)

    def y_range(self, number_CSV, csv_list_path, csv_list, plot_y_range_positive, plot_y_range_negative):
        for j in range(0, number_CSV, 1):
            hysteresis_temp_df = pd.read_csv(csv_list_path + csv_list[j], header=None, engine='c')
            y = hysteresis_temp_df.iloc[:, 1]
            YMax = y.max(axis=0)
            YMin = y.min(axis=0)
            if YMax > plot_y_range_positive:
                plot_y_range_positive = YMax
            if YMin < plot_y_range_negative:
                plot_y_range_negative = YMin

        if abs(plot_y_range_positive) > abs(plot_y_range_negative):
            hyst_lim = abs(plot_y_range_positive)
        else:
            hyst_lim = abs(plot_y_range_negative)
        return hyst_lim

    def process_individual_hysteresis(self, number_CSV, csv_list, csv_list_path, hyst_lim, folder):
        for j in range(0, number_CSV, 1):
            hysteresis_temp_df = pd.read_csv(csv_list_path + csv_list[j], header=0, engine='c')
            thermal = hysteresis_temp_df.iloc[:, 0]
            cur_temp = round(thermal.mean(), 1)
            cur_temp_file = str(cur_temp).replace(".", "_")
            # cur_temp_file = cur_temp_file[:cur_temp_file.rfind('_')] if cur_temp_file[-2:] == '_0' else cur_temp_file
            if cur_temp_file.endswith('_0'):
                cur_temp_file = cur_temp_file[:cur_temp_file.rfind('_')]

            # Case 2: ends with "_9" → parse the part before "_" as int, then +1
            elif cur_temp_file.endswith('_9'):
                head = cur_temp_file[:cur_temp_file.rfind('_')]
                try:
                    cur_temp_file = str(int(head) + 1)
                except ValueError:
                    # Fallback if the part before '_' isn't a clean integer
                    # (optional) leave unchanged or handle/log as needed
                    pass
            x = hysteresis_temp_df.iloc[:, 0]
            y = hysteresis_temp_df.iloc[:, 1]

            fig, ax_hys = plt.subplots()
            ax_hys.scatter(x, y, color='black', s=0.5)
            ax_hys.set_xlabel('Magnetic Field (Oe)', fontsize=14)
            ax_hys.set_ylabel('Magnetic Moment (emu)', fontsize=14)
            ax_hys.set_ylim(bottom=hyst_lim * -1.05, top=hyst_lim * 1.05)
            plt.title('RAW {} K Hysteresis Loop'.format(cur_temp), pad=10, wrap=True, fontsize=14)
            plt.tight_layout()
            plt.savefig(folder + "/Raw_{}K_Hysteresis.png".format(cur_temp_file))
            plt.close()

    def process_data(self):
        if not hasattr(self, 'folder_selected') or self.save_folder_path is None:
            return

        # Update progress
        self.progress_bar.setValue(5)
        QApplication.processEvents()

        self.ProcessedRAW = self.save_folder_path + '/Processed_Graph'
        isExist = os.path.exists(self.ProcessedRAW)
        if not isExist:
            os.makedirs(self.ProcessedRAW)

        self.ProcCal = self.save_folder_path + '/Processed_Result'
        isExist = os.path.exists(self.ProcCal)
        if not isExist:
            os.makedirs(self.ProcCal)

        # Get files from tree view instead of scanning directory
        csv_list = []
        root = self.file_tree.invisibleRootItem()
        child_count = root.childCount()

        for i in range(child_count):
            item = root.child(i)
            file_path = item.toolTip(0)
            csv_list.append((os.path.basename(file_path), file_path))

        number_CSV = len(csv_list)

        if number_CSV == 0:
            QMessageBox.warning(self, "Error",
                                "No CSV files found in tree view.\n\n" +
                                "Please ensure your CSV files:\n" +
                                "1. Are named like: 10K.csv, 50K.csv, 300K.csv\n" +
                                "2. Have columns: Temperature (K), Magnetic Field (Oe), Moment (emu)\n" +
                                "3. Have headers in the first row")
            return

        self.progress_bar.setValue(10)
        QApplication.processEvents()

        plot_y_range_positive = 0
        plot_y_range_negative = 0

        for j in range(0, number_CSV, 1):
            file_name, file_path = csv_list[j]
            if file_name.lower() != 'zfc.csv' and file_name.lower() != 'fc.csv':
                try:
                    hysteresis_temp_df = pd.read_csv(file_path, header=0, engine='c')
                    YMax = hysteresis_temp_df['Moment (emu)'].max(axis=0)
                    YMin = hysteresis_temp_df['Moment (emu)'].min(axis=0)

                    if YMax > plot_y_range_positive:
                        plot_y_range_positive = YMax
                    if YMin < plot_y_range_negative:
                        plot_y_range_negative = YMin
                except Exception as e:
                    print(f"Error processing {file_name}: {e}")
                    continue

        if abs(plot_y_range_positive) > abs(plot_y_range_negative):
            hyst_lim = abs(plot_y_range_positive)
        else:
            hyst_lim = abs(plot_y_range_negative)

        # area_df = pd.DataFrame(columns=['Temperature', 'Area'])
        area_list = []
        # Ms_df = pd.DataFrame(columns=['Temperature', 'Saturation Field', 'Upper', 'Lower'])
        ms_list = []
        # Coercivity_df = pd.DataFrame(columns=['Temperature', 'Coercivity'])
        coercivity_list = []
        # eb_df = pd.DataFrame(columns=['Temperature', 'Exchange Bias'])
        eb_list = []

        self.progress_bar.setValue(15)
        QApplication.processEvents()

        for j in range(0, number_CSV, 1):
            # Update progress bar for each file
            progress = 15 + int((j / number_CSV) * 60)  # 15% to 75%
            self.progress_bar.setValue(progress)
            self.progress_label.setText(f"Processing file {j + 1}/{number_CSV}...")
            QApplication.processEvents()

            file_name, file_path = csv_list[j]
            self.hysteresis_folder = self.ProcessedRAW + '/Hysteresis'
            isExist = os.path.exists(self.hysteresis_folder)
            if not isExist:
                os.makedirs(self.hysteresis_folder)

            hysteresis_temp_df = pd.read_csv(file_path, header=0, engine='c')
            temperature_column = hysteresis_temp_df.iloc[:, 0]

            first_item = temperature_column.iloc[0]  # First item
            last_item = temperature_column.iloc[-1]  # Last item
            temperature_difference = first_item - last_item
            if file_name == 'zfc.csv' or file_name == 'fc.csv' or abs(temperature_difference) > 5:
                ZFC_FC = pd.read_csv(file_path, header=0, engine='c')
                x = ZFC_FC.iloc[:, 0]
                field = ZFC_FC.iloc[:, 1]
                field_mean = int(round(field.mean()))
                y = ZFC_FC.iloc[:, 2]
                plt.rc('xtick', labelsize=13)
                plt.rc('ytick', labelsize=13)
                fig, ax = plt.subplots()
                ax.scatter(x, y, color='black', s=0.5)
                ax.set_xlabel('Magnetic Field (Oe)', fontsize=14)
                ax.set_ylabel('Magnetic Moment (emu)', fontsize=14)
                ax.set_ylim(bottom=hyst_lim * -1.05, top=hyst_lim * 1.05)
                ax.text(0.95, 0.95, f'{field_mean} Oe', transform=ax.transAxes,
                        fontsize=12, ha='right', va='top')
                if file_name == 'zfc.csv' or temperature_difference < 0 and field_mean == 0:
                    plt.title('Zero Field Cooled', pad=10, wrap=True, fontsize=14)

                    plt.tight_layout()
                    plt.savefig(self.hysteresis_folder + "/zfc.png")
                elif file_name == 'fcc.csv' or temperature_difference > 0 and field_mean > 0:
                    plt.title('Field Cooled Cooling', pad=10, wrap=True, fontsize=14)
                    plt.tight_layout()
                    plt.savefig(self.hysteresis_folder + "/fcc.png")
                elif file_name == 'fcw.csv' or temperature_difference < 0 and field_mean > 0:
                    plt.title('Field Cooled Warming', pad=10, wrap=True, fontsize=14)
                    plt.tight_layout()
                    plt.savefig(self.hysteresis_folder + "/fcw.png")
                else:
                    plt.title('FC or ZFC', pad=10, wrap=True, fontsize=14)
                    plt.tight_layout()
                    plt.savefig(self.hysteresis_folder + "/fc_zfc_unknown.png")
                plt.close()
            else:
                try:
                    cur_temp = round(temperature_column.mean(), 1)
                    cur_temp_file = str(cur_temp).replace(".", "_")

                    if cur_temp_file.endswith('_0'):
                        cur_temp_file = cur_temp_file[:cur_temp_file.rfind('_')]

                    # Case 2: ends with "_9" → parse the part before "_" as int, then +1
                    elif cur_temp_file.endswith('_9'):
                        head = cur_temp_file[:cur_temp_file.rfind('_')]
                        try:
                            cur_temp_file = str(int(head) + 1)
                        except ValueError:
                            # Fallback if the part before '_' isn't a clean integer
                            # (optional) leave unchanged or handle/log as needed
                            pass


                    x = hysteresis_temp_df.iloc[:, 1]
                    y = hysteresis_temp_df.iloc[:, 2]

                    self.raw_data_temp_folder = self.ProcCal + '/raw_data_with_temperature'
                    isExist = os.path.exists(self.raw_data_temp_folder)
                    if not isExist:
                        os.makedirs(self.raw_data_temp_folder)
                    raw_temp_df = pd.concat([x, y],
                                         ignore_index=True, axis=1)
                    raw_temp_df.to_csv(self.raw_data_temp_folder + '/{}K.csv'.format(cur_temp_file),
                                    index=False, header=False)

                    # This section is for the area calculation
                    hysteresis_temp_df_area = hysteresis_temp_df.dropna()
                    x_drop = hysteresis_temp_df_area.iloc[:, 1]
                    y_drop = hysteresis_temp_df_area.iloc[:, 2]
                    area = np.trapz(y_drop, x_drop)
                    area_list.append({'Temperature': cur_temp, 'Area': abs(area)})
                    # area_temp = pd.DataFrame({'Temperature': [cur_temp], 'Area': [abs(area)]})
                    # area_df = pd.concat([area_df, area_temp], ignore_index=True)
                    plt.plot(x_drop, y_drop)
                    plt.fill_between(x_drop, y_drop, alpha=0.3)
                    plt.legend()
                    plt.title("Trapezoidal Area")
                    plt.xlabel('Temperature (K)')
                    plt.ylabel('Magnetic Moment (emu)')
                    plt.tight_layout()
                    self.Area_folder = self.ProcessedRAW + '/Area'
                    isExist = os.path.exists(self.Area_folder)
                    if not isExist:
                        os.makedirs(self.Area_folder)
                    plt.savefig(self.Area_folder + "/{}K_Area.png".format(cur_temp_file))
                    plt.close()

                    # This section is for thermal difference plot
                    thermal_difference = 100 * (temperature_column - cur_temp) / cur_temp
                    plt.rc('xtick', labelsize=13)
                    plt.rc('ytick', labelsize=13)
                    fig, ax_thermal = plt.subplots()
                    ax_thermal.scatter(x, thermal_difference, color='black', s=0.5)
                    ax_thermal.set_xlabel('Magnetic Field (Oe)', fontsize=14)
                    ax_thermal.set_ylabel('Temperature Fluctuation (%)', fontsize=14)
                    plt.title('{} K Hysteresis Loop Temperature'.format(cur_temp), pad=10, wrap=True, fontsize=14)
                    plt.tight_layout()
                    self.thermal_difference_folder = self.ProcessedRAW + '/Thermal_difference'
                    isExist = os.path.exists(self.thermal_difference_folder)
                    if not isExist:
                        os.makedirs(self.thermal_difference_folder)
                    plt.savefig(self.thermal_difference_folder + "/{}K_Thermal_Difference.png".format(cur_temp_file))
                    plt.close()
                    thermal_difference_x = pd.Series(x)
                    thermal_difference_y = pd.Series(thermal_difference)
                    thermal_difference_df = pd.concat([thermal_difference_x, thermal_difference_y],
                                         ignore_index=True, axis=1)
                    self.thermal_difference_folder = self.ProcCal + '/Thermal_Difference'
                    isExist = os.path.exists(self.thermal_difference_folder)
                    if not isExist:
                        os.makedirs(self.thermal_difference_folder)
                    thermal_difference_df.to_csv(self.thermal_difference_folder + '/{}K_thermal_difference.csv'.format(cur_temp_file),
                                    index=False, header=False)

                    # This section is plotting the RAW hysteresis data
                    fig, ax_hys = plt.subplots()
                    ax_hys.scatter(x, y, color='black', s=0.5)
                    ax_hys.set_xlabel('Magnetic Field (Oe)', fontsize=14)
                    ax_hys.set_ylabel('Magnetic Moment (emu)', fontsize=14)
                    ax_hys.set_ylim(bottom=hyst_lim * -1.05, top=hyst_lim * 1.05)
                    plt.title('{} K Hysteresis Loop'.format(cur_temp), pad=10, wrap=True, fontsize=14)
                    plt.tight_layout()
                    plt.savefig(self.hysteresis_folder + "/{}K_RAW_Hysteresis.png".format(cur_temp_file))
                    plt.close()

                    # This section split the raw data into upper half and lower half data
                    self.split_folder = self.ProcCal + '/Split_folder'
                    isExist = os.path.exists(self.split_folder)
                    if not isExist:
                        os.makedirs(self.split_folder)

                    index = hysteresis_temp_df.iloc[:, 1].idxmin(axis=0)
                    list1 = hysteresis_temp_df.iloc[:index].dropna()
                    list2 = hysteresis_temp_df.iloc[index:].dropna()
                    list1_x = list1["Magnetic Field (Oe)"].values
                    list1_y = list1["Moment (emu)"].values
                    list2_x = list2["Magnetic Field (Oe)"].values
                    list2_y = list2["Moment (emu)"].values
                    fig, axs = plt.subplots()
                    plt.title("{}K Split Hysteresis".format(cur_temp), pad=10, wrap=True, fontsize=14)
                    axs.scatter(list1_x, list1_y, s=10)
                    axs.scatter(list2_x, list2_y, s=10, alpha=0.1)
                    ax_hys.set_xlabel('Magnetic Field (Oe)', fontsize=14)
                    ax_hys.set_ylabel('Magnetic Moment (emu)', fontsize=14)
                    plt.tight_layout()
                    plt.savefig(self.split_folder + "/{}K_spliting.png".format(cur_temp_file))
                    plt.close()

                    list1_x_concat = pd.Series(list1_x)
                    list1_y_concat = pd.Series(list1_y)
                    list2_x_concat = pd.Series(list2_x)
                    list2_y_concat = pd.Series(list2_y)
                    spilt_df = pd.concat([list1_x_concat, list1_y_concat, list2_x_concat, list2_y_concat],
                                         ignore_index=True, axis=1)
                    self.split_raw = self.ProcCal + '/Raw_Data_Spilt'
                    isExist = os.path.exists(self.split_raw)
                    if not isExist:
                        os.makedirs(self.split_raw)
                    spilt_df.to_csv(self.split_raw + '/{}K_spliting_raw_data.csv'.format(cur_temp_file),
                                    index=False, header=False)

                    # This section is to find the linear background
                    def find_slope(x: list, y: list, portion:str=None):
                        pm_dm_result = extract_pm_dm_slope(
                            x, y,
                            method='linear_saturation',
                            saturation_threshold=0.70
                        )

                        # Check quality
                        if pm_dm_result['r_squared'] < 0.8:
                            print(f"    Warning: Low R² = {pm_dm_result['r_squared']:.4f} in {portion} using linear saturation method")

                            # Try symmetric method as backup
                            pm_dm_result_sym = extract_pm_dm_slope(
                                x, y,
                                method='symmetric',
                                saturation_threshold=0.70
                            )
                            if pm_dm_result_sym['r_squared'] > pm_dm_result['r_squared']:
                                pm_dm_result = pm_dm_result_sym
                                print(f"    Using symmetric method instead in {portion}")
                                # Try high field method as backup
                            if pm_dm_result['r_squared'] < 0.84 and pm_dm_result_sym['r_squared'] < 0.85:
                                pm_dm_result_high_field = extract_pm_dm_slope(
                                    x, y,
                                    method='high_field_only',
                                    saturation_threshold=0.70
                                )
                                pm_dm_result = pm_dm_result_high_field
                                print(f"    Using high field method instead in {portion}")



                        chi_total = pm_dm_result['chi_total']
                        pm_dm_method = pm_dm_result['method']

                        # print(f"    χ_total (PM+DM) = {chi_total:.6e} emu/Oe")
                        # print(f"    R² = {pm_dm_result.get('r_squared', 0):.4f}")

                        # Determine type
                        if chi_total > 1e-10:
                            chi_type = 'Paramagnetic'
                            # print(f"    → Paramagnetic contribution (positive slope)")
                        elif chi_total < -1e-10:
                            chi_type = 'Diamagnetic'
                            # print(f"    → Diamagnetic contribution (negative slope)")
                        else:
                            chi_type = 'Negligible'
                            # print(f"    → Negligible PM/DM contribution")

                        # Save PM/DM diagnostics
                        ### NEW ### Create PM/DM folder
                        self.fit_folder = self.ProcessedRAW + '/PM_DM_Analysis'
                        isExist = os.path.exists(self.fit_folder)
                        if not isExist:
                            os.makedirs(self.fit_folder)
                        self.pm_dm_folder = self.fit_folder + '/PM_DM_Analysis'
                        os.makedirs(self.pm_dm_folder, exist_ok=True)

                        fig_pm_dm = visualize_pm_dm_extraction(x, y, pm_dm_result)
                        fig_pm_dm.savefig(f"{self.pm_dm_folder}/{cur_temp_file}K_PM_DM_extraction_{portion}.png", dpi=150)
                        plt.close(fig_pm_dm)
                        return pm_dm_result

                    try:
                        # Try linear saturation method first (best)
                        pm_dm_result_upper = find_slope(list1_x, list1_y, portion="Upper")
                        pm_dm_result_lower = find_slope(list2_x, list2_y, portion='Lower')
                        slope_upper = pm_dm_result_upper['chi_total']
                        slope_lower = pm_dm_result_lower['chi_total']
                        if slope_lower and slope_upper:
                            slope = (slope_lower + slope_upper) / 2
                        elif not slope_lower:
                            slope = slope_upper
                        else:
                            slope = slope_lower
                        slope_upper_x_concat = pd.Series(list1_x)
                        slope_y_concat = pd.Series(list1_x * slope)
                        slope_df = pd.concat([slope_upper_x_concat, slope_y_concat],
                                             ignore_index=True, axis=1)
                        self.slope_folder = self.ProcCal + '/Slope'
                        isExist = os.path.exists(self.slope_folder)
                        if not isExist:
                            os.makedirs(self.slope_folder)
                        slope_df.to_csv(self.slope_folder + '/{}K_slope_data.csv'.format(cur_temp_file),
                                        index=False, header=False)

                        # This section is to remove the linear background
                        M_corrected_upper = remove_pm_dm_contribution(list1_x, list1_y, slope,
                                                                remove_offset=False)
                        M_corrected_lower = remove_pm_dm_contribution(list2_x, list2_y, slope,
                                                                remove_offset=False)

                        y_slope_removal_upper_concat = pd.Series(M_corrected_upper)
                        y_slope_removal_lower_concat = pd.Series(M_corrected_lower)
                        slope_removal_spilt_data_df = pd.concat([list1_x_concat, y_slope_removal_upper_concat,
                                              list2_x_concat, y_slope_removal_lower_concat],
                                             ignore_index=True, axis=1)
                        self.slope_removal_spilt_folder = self.ProcCal + '/Slope_Removal_Spilt'
                        isExist = os.path.exists(self.slope_removal_spilt_folder)
                        if not isExist:
                            os.makedirs(self.slope_removal_spilt_folder)
                        slope_removal_spilt_data_df.to_csv(self.slope_removal_spilt_folder + '/{}K_slope_removal_spilt_data.csv'.format(cur_temp_file),
                                        index=False, header=False)

                        x_all_slope_removed_processed = pd.concat([list1_x_concat, list2_x_concat], ignore_index=True)
                        y_all_slope_removed_processed = pd.concat([y_slope_removal_upper_concat, y_slope_removal_lower_concat], ignore_index=True)

                        slope_removal_data_df = pd.concat([x_all_slope_removed_processed, y_all_slope_removed_processed], axis=1)
                        self.slope_removal_older = self.ProcCal + '/Slope_Removal'
                        isExist = os.path.exists(self.slope_removal_older)
                        if not isExist:
                            os.makedirs(self.slope_removal_older)
                        slope_removal_data_df.to_csv(
                            self.slope_removal_older + '/{}K_slope_removal_data.csv'.format(cur_temp_file),
                            index=False, header=False)

                    except Exception as e:
                        print(f"    Error extracting PM/DM: {e}")
                        chi_total = 0
                        chi_type = 'Failed'
                        pm_dm_method = 'None'
                        pm_dm_result = {'chi_total': 0, 'r_squared': 0, 'n_points': 0}

                    Ms_result_upper = extract_saturation_magnetization(list1_x_concat, y_slope_removal_upper_concat, n_points=5)
                    Ms_result_lower = extract_saturation_magnetization(list2_x_concat, y_slope_removal_lower_concat,
                                                                 n_points=5)
                    ms_upper = (Ms_result_upper['Ms_upper'] + Ms_result_lower['Ms_upper']) / 2
                    ms_lower = (Ms_result_upper['Ms_lower'] + Ms_result_lower['Ms_lower']) / 2
                    ms_avg = (Ms_result_upper['Ms_avg'] + Ms_result_lower['Ms_avg']) / 2
                    ms_vertical_offset = (Ms_result_upper['vertical_offset'] + Ms_result_lower['vertical_offset']) / 2

                    # Step 2: Correct vertical offset if needed
                    if abs(ms_vertical_offset) > 1e-10:
                        M_upper_corrected = correct_vertical_offset(y_slope_removal_upper_concat, ms_vertical_offset)
                        M_lower_corrected = correct_vertical_offset(y_slope_removal_lower_concat, ms_vertical_offset)
                        M_all_corrected = correct_vertical_offset(y_all_slope_removed_processed, ms_vertical_offset)
                    else:
                        M_upper_corrected = y_slope_removal_upper_concat.copy()
                        M_lower_corrected = y_slope_removal_lower_concat.copy()
                        M_all_corrected = y_all_slope_removed_processed.copy()
                    # Re-calculate Ms after correction

                    ms_corrected_spilt_data_df = pd.concat([list1_x_concat, M_upper_corrected,
                                                             list2_x_concat, M_lower_corrected],
                                                            ignore_index=True, axis=1)
                    self.ms_corrected_spilt_folder = self.ProcCal + '/Ms_Corrected_Spilt'
                    isExist = os.path.exists(self.ms_corrected_spilt_folder)
                    if not isExist:
                        os.makedirs(self.ms_corrected_spilt_folder)
                    ms_corrected_spilt_data_df.to_csv(
                        self.ms_corrected_spilt_folder + '/{}K_ms_corrected_spilt_data.csv'.format(cur_temp_file),
                        index=False, header=False)

                    ms_corrected_data_df = pd.concat([x_all_slope_removed_processed, M_all_corrected],
                                                           ignore_index=True, axis=1)
                    self.ms_corrected_folder = self.ProcCal + '/Ms_Corrected'
                    isExist = os.path.exists(self.ms_corrected_folder)
                    if not isExist:
                        os.makedirs(self.ms_corrected_folder)
                    ms_corrected_data_df.to_csv(
                        self.ms_corrected_folder + '/{}K_ms_corrected_data.csv'.format(cur_temp_file),
                        index=False, header=False)

                    # Ms_temp = pd.DataFrame(
                    #     {'Temperature': [cur_temp], 'Saturation Field': [ms_avg], 'Upper': [ms_upper],
                    #      'Lower': [ms_lower]})
                    # Ms_df = pd.concat([Ms_df, Ms_temp], ignore_index=True)
                    ms_list.append({'Temperature': cur_temp, 'Saturation Field': ms_avg,
                                    'Upper': ms_upper, 'Lower': ms_lower})

                    upper_fit = extract_coercivity_method1_fit(list1_x_concat, M_upper_corrected, include_slope=False)

                    self.fit_summary_folder = self.ProcCal + '/Fit_Summary'
                    isExist = os.path.exists(self.fit_summary_folder)
                    if not isExist:
                        os.makedirs(self.fit_summary_folder)

                    self.fitted_curve_folder = self.ProcCal + '/Fitted_Curve'
                    isExist = os.path.exists(self.fitted_curve_folder)
                    if not isExist:
                        os.makedirs(self.fitted_curve_folder)

                    if upper_fit['success']:
                        upper_fit_result = pd.DataFrame(upper_fit)
                        upper_fit_result.to_csv(
                            self.fit_summary_folder + '/{}K_ms_corrected_data_fitting_results_upper_curve.csv'.format(cur_temp_file),
                            index=False, header=False)

                        upper_fit_concat = pd.concat([list1_x_concat, upper_fit['M_fit']],
                                                         ignore_index=True, axis=1)

                        upper_fit_concat.to_csv(
                            self.fitted_curve_folder + '/{}K_ms_corrected_data_fitting_curve_upper_curve.csv'.format(cur_temp_file),
                            index=False, header=False)
                    else:
                        print(f"  ✗ Upper fit failed: {upper_fit.get('error', 'Unknown error')}")


                    # Step 4: Method 1 - Fit lower branch
                    lower_fit = extract_coercivity_method1_fit(list2_x_concat, M_lower_corrected, include_slope=False)

                    if lower_fit['success']:
                        lower_fit_result = pd.DataFrame(lower_fit)
                        lower_fit_result.to_csv(
                            self.fit_summary_folder + '/{}K_ms_corrected_data_fitting_results_lower_curve.csv'.format(
                                cur_temp_file),
                            index=False, header=False)
                        low_fit_concat = pd.concat([list2_x_concat, lower_fit['M_fit']],
                                                     ignore_index=True, axis=1)
                        low_fit_concat.to_csv(
                            self.fitted_curve_folder + '/{}K_ms_corrected_data_fitting_curve_lower_curve.csv'.format(cur_temp_file),
                            index=False, header=False)
                    else:
                        print(f"  ✗ Lower fit failed: {lower_fit.get('error', 'Unknown error')}")

                    # Step 5: Method 2 - Data crossing
                    method2_result = extract_coercivity_method2_data(list1_x, M_upper_corrected,
                                                                     list2_x, M_lower_corrected)
                    if method2_result['success']:
                        method2_fit_result = pd.DataFrame(method2_result)
                        method2_fit_result.to_csv(
                            self.fit_summary_folder + '/{}K_direct_result.csv'.format(
                                cur_temp_file),
                            index=False, header=False)
                    else:
                        print(f"  ✗ Could not find crossings: {method2_result.get('error', 'Unknown')}")

                    # Step 6: Determine final values
                    Hc_values = []
                    Heb_values = []

                    # Hc from upper branch fit
                    if upper_fit['success'] and upper_fit['r_squared'] > 0.90 and lower_fit['success'] and lower_fit['r_squared'] > 0.90:
                        Hc_left = upper_fit['c']
                        Hc_right = lower_fit['c']
                        Hc = abs(Hc_left + Hc_right)
                        Hc_values.append(Hc)
                        Heb = Hc_left + Hc_right
                        Heb_values.append(Heb)

                    # Hc from Method 2 (data crossings)
                    if method2_result['success']:
                        Hc_values.append(method2_result['Hc'])
                        Heb_values.append(method2_result['Heb'])

                    if len(Hc_values) > 0:
                        Hc_final = np.mean(Hc_values)
                    else:
                        Hc_final = None

                    if len(Heb_values) > 0:
                        Heb_final = np.mean(Heb_values)
                    else:
                        Heb_final = None

                    # coercivity_temp = pd.DataFrame({'Temperature': [cur_temp],
                    #                                     'Coercivity': [Hc_final]})
                    # Coercivity_df = pd.concat([Coercivity_df, coercivity_temp], ignore_index=True)
                    coercivity_list.append({'Temperature': cur_temp, 'Coercivity': Hc_final})

                    # eb_temp = pd.DataFrame({'Temperature': [cur_temp],
                    #                                 'Exchange Bias': [Heb_final]})
                    # eb_df = pd.concat([eb_df, eb_temp], ignore_index=True)
                    eb_list.append({'Temperature': cur_temp, 'Exchange Bias': Heb_final})

                    eb_all_corrected = x_all_slope_removed_processed - Heb_final
                    eb_ms_corrected_data_df = pd.concat([eb_all_corrected, M_all_corrected],
                                                     ignore_index=True, axis=1)

                    self.ms_eb_corrected_folder = self.ProcCal + '/Ms_Eb_Corrected'
                    isExist = os.path.exists(self.ms_eb_corrected_folder)
                    if not isExist:
                        os.makedirs(self.ms_eb_corrected_folder)
                    eb_ms_corrected_data_df.to_csv(
                        self.ms_eb_corrected_folder + '/{}K_eb_ms_corrected_data.csv'.format(cur_temp_file),
                        index=False, header=False)

                    eb_corrected_data_df = pd.concat([eb_all_corrected, y_all_slope_removed_processed],
                                                        ignore_index=True, axis=1)
                    self.eb_corrected_folder = self.ProcCal + '/Eb_Corrected'
                    isExist = os.path.exists(self.eb_corrected_folder)
                    if not isExist:
                        os.makedirs(self.eb_corrected_folder)
                    eb_corrected_data_df.to_csv(
                        self.eb_corrected_folder + '/{}K_eb_corrected_data.csv'.format(cur_temp_file),
                        index=False, header=False)


                except Exception as e:
                    tb_str = traceback.format_exc()
                    print(f"Error processing temperature {file_name}: {e} {tb_str}")
                    continue

        # Sort and save summary data
        self.progress_bar.setValue(80)
        self.progress_label.setText("Generating summary plots...")
        QApplication.processEvents()

        area_df = pd.DataFrame(area_list) if area_list else pd.DataFrame(columns=['Temperature', 'Area'])
        area_df = area_df.sort_values('Temperature')
        x_temp = area_df.iloc[:, 0]
        y_area = area_df.iloc[:, 1]

        plt.rc('xtick', labelsize=13)
        plt.rc('ytick', labelsize=13)
        fig, ax = plt.subplots()
        ax.plot(x_temp, y_area, color='black', marker='s', linewidth=2, markersize=5)
        ax.set_xlabel('Temperature (K)', fontsize=14)
        ax.set_ylabel('Area', fontsize=14)
        plt.tight_layout()
        plt.savefig(self.ProcessedRAW + "/Area_Relation.png")
        area_df.to_csv(self.ProcCal + '/Area_Relation.csv', index=False)
        plt.close()

        ms_df = pd.DataFrame(ms_list) if ms_list else pd.DataFrame(columns=['Temperature', 'Saturation Field', 'Upper', 'Lower'])
        ms_df = ms_df.sort_values('Temperature')
        x_temp = ms_df.iloc[:, 0]
        y_ms = ms_df.iloc[:, 1]
        y_ms_upper = ms_df.iloc[:, 2]
        y_ms_lower = ms_df.iloc[:, 3]

        plt.rc('xtick', labelsize=13)
        plt.rc('ytick', labelsize=13)
        fig, ax = plt.subplots()
        ax.plot(x_temp, y_ms, color='black', marker='s', linewidth=2, markersize=5)
        ax.set_xlabel('Temperature (K)', fontsize=14)
        ax.set_ylabel('Saturation Field Ms', fontsize=14)
        plt.tight_layout()
        plt.savefig(self.ProcessedRAW + "/Saturation_Field.png")
        ms_df.to_csv(self.ProcCal + '/Saturation_Field.csv', index=False)
        plt.close()

        plt.rc('xtick', labelsize=13)
        plt.rc('ytick', labelsize=13)
        fig, ax = plt.subplots()
        ax.plot(x_temp, y_ms_upper, color='#922B21', marker='s', linewidth=2, markersize=5, label="upper")
        ax.plot(x_temp, y_ms_lower, color='#212F3D', marker='s', linewidth=2, markersize=5, label="lower")
        ax.set_xlabel('Temperature (K)', fontsize=14)
        ax.set_ylabel('Saturation Field Ms', fontsize=14)
        plt.tight_layout()
        plt.legend()
        plt.savefig(self.ProcessedRAW + "/Saturation_Field_Split.png")
        plt.close()

        coercivity_df = pd.DataFrame(coercivity_list) if coercivity_list else pd.DataFrame(
            columns=['Temperature', 'Coercivity'])
        coercivity_df = coercivity_df.sort_values('Temperature')
        x_temp = coercivity_df.iloc[:, 0]
        y_coerc = coercivity_df.iloc[:, 1]

        plt.rc('xtick', labelsize=13)
        plt.rc('ytick', labelsize=13)
        fig, ax = plt.subplots()
        ax.plot(x_temp, y_coerc, color='black', marker='s', linewidth=2, markersize=5)
        ax.set_xlabel('Temperature (K)', fontsize=14)
        ax.set_ylabel('Coercivity', fontsize=14)
        plt.tight_layout()
        plt.savefig(self.ProcessedRAW + "/Coercivity.png")
        coercivity_df.to_csv(self.ProcCal + '/Coercivity.csv', index=False)
        plt.close()

        eb_df = pd.DataFrame(eb_list) if eb_list else pd.DataFrame(
            columns=['Temperature', 'Exchange Bias'])
        eb_df = eb_df.sort_values('Temperature')
        x_temp = eb_df.iloc[:, 0]
        y_rb = eb_df.iloc[:, 1]

        plt.rc('xtick', labelsize=13)
        plt.rc('ytick', labelsize=13)
        fig, ax = plt.subplots()
        ax.plot(x_temp, y_rb, color='black', marker='s', linewidth=2, markersize=5)
        ax.set_xlabel('Temperature (K)', fontsize=14)
        ax.set_ylabel('Exchange Bias', fontsize=14)
        plt.tight_layout()
        plt.savefig(self.ProcessedRAW + "/Exchange_bias.png")
        eb_df.to_csv(self.ProcCal + '/Exchange_bias.csv', index=False)
        plt.close()

        # Process individual hysteresis plots
        self.progress_bar.setValue(85)
        self.progress_label.setText("Processing individual hysteresis plots...")
        QApplication.processEvents()

        y_range_positive = 0
        y_range_negative = 0
        data_csv_list_path = self.ProcCal + '/Slope_Removal/'

        if os.path.exists(data_csv_list_path):
            data_csv_list = os.listdir(data_csv_list_path)
            data_number_CSV = len(data_csv_list)
            self.Final_Hysteresis = self.ProcessedRAW + '/Final_Hysteresis'
            isExist = os.path.exists(self.Final_Hysteresis)
            if not isExist:
                os.makedirs(self.Final_Hysteresis)

            data_hyst_lim = self.y_range(data_number_CSV, data_csv_list_path, data_csv_list,
                                         y_range_positive, y_range_negative)
            self.process_individual_hysteresis(data_number_CSV, data_csv_list, data_csv_list_path,
                                        data_hyst_lim, self.Final_Hysteresis)

        # y_range_positive = 0
        # y_range_negative = 0
        # data_csv_list_path = self.save_folder_path + '/'
        #
        # if os.path.exists(data_csv_list_path):
        #     data_csv_list = os.listdir(data_csv_list_path)
        #     data_number_CSV = len(data_csv_list)
        #     self.Final_Hysteresis_RAW = self.ProcessedRAW + '/Raw_Hysteresis'
        #     isExist = os.path.exists(self.Final_Hysteresis_RAW)
        #     if not isExist:
        #         os.makedirs(self.Final_Hysteresis_RAW)
        #
        #     data_hyst_lim = self.y_range(data_number_CSV, data_csv_list_path, data_csv_list,
        #                                  y_range_positive, y_range_negative)
        #     self.process_individual_hysteresis(data_number_CSV, data_csv_list, data_csv_list_path,
        #                                 data_hyst_lim, self.Final_Hysteresis_RAW)

        # Generate PowerPoint
        self.progress_bar.setValue(95)
        self.progress_label.setText("Generating PowerPoint presentations...")
        QApplication.processEvents()

        self.to_ppt()

        # Load summary data for interactive display
        self.progress_bar.setValue(98)
        self.progress_label.setText("Loading summary data...")
        QApplication.processEvents()

        try:
            area_path = os.path.join(self.ProcCal, 'Area_Relation.csv')
            if os.path.exists(area_path):
                self.area_df = pd.read_csv(area_path)

            ms_path = os.path.join(self.ProcCal, 'Saturation_Field.csv')
            if os.path.exists(ms_path):
                self.Ms_df = pd.read_csv(ms_path)

            coercivity_path = os.path.join(self.ProcCal, 'Coercivity.csv')
            if os.path.exists(coercivity_path):
                self.Coercivity_df = pd.read_csv(coercivity_path)

            eb_path = os.path.join(self.ProcCal, 'Exchange_bias.csv')
            if os.path.exists(coercivity_path):
                self.eb_df = pd.read_csv(eb_path)

            # Update summary plot if any checkboxes are checked
            self.update_summary_plot()

        except Exception as e:
            print(f"Error loading summary data: {e}")

        # Complete
        self.progress_bar.setValue(100)
        self.progress_label.setText("Processing complete!")
        QApplication.processEvents()

        # Enable file selection after processing
        self.is_processed = True
        self.file_tree.setSelectionMode(QTreeWidget.SelectionMode.SingleSelection)

        # Show completion message
        QMessageBox.information(self, "Processing Complete",
                                f"Successfully processed {number_CSV} files.\n\n" +
                                "Results saved to:\n" +
                                f"- {self.ProcessedRAW}\n" +
                                f"- {self.ProcCal}\n\n" +
                                "You can now select files from the tree to view detailed plots.")

    def rstpage(self):
        try:
            # Don't try to reset if widgets don't exist yet
            if hasattr(self, 'file_tree'):
                self.file_tree.clear()
        except Exception as e:
            print(f"Reset page warning: {e}")
            pass

    def to_ppt(self):
        """Generate PowerPoint presentations with hysteresis plots"""
        try:
            if os.path.exists(self.ProcCal + '/slope_removal.pptx'):
                prs = Presentation(self.ProcCal + '/slope_removal.pptx')
                prs.slide_width = Inches(13.33)
                prs.slide_height = Inches(7.5)
            else:
                prs = Presentation()
                prs.slide_width = Inches(13.33)
                prs.slide_height = Inches(7.5)
                prs.save(self.ProcCal + '/slope_removal.pptx')
                prs = Presentation(self.ProcCal + '/slope_removal.pptx')
                prs.slide_width = Inches(13.33)
                prs.slide_height = Inches(7.5)

            # if os.path.exists(self.ProcCal + '/raw_hysteresis.pptx'):
            #     prs_RAW = Presentation(self.ProcCal + '/raw_hysteresis.pptx')
            #     prs_RAW.slide_width = Inches(13.33)
            #     prs_RAW.slide_height = Inches(7.5)
            # else:
            #     prs_RAW = Presentation()
            #     prs_RAW.slide_width = Inches(13.33)
            #     prs_RAW.slide_height = Inches(7.5)
            #     prs_RAW.save(self.ProcCal + '/raw_hysteresis.pptx')
            #     prs_RAW = Presentation(self.ProcCal + '/raw_hysteresis.pptx')
            #     prs_RAW.slide_width = Inches(13.33)
            #     prs_RAW.slide_height = Inches(7.5)

            dat_list_path = self.folder_selected + '*.DAT'
            dat_list = glob.glob(dat_list_path)

            if len(dat_list) > 0:
                file_name = dat_list[0]
                for i in range(len(file_name) - 1, 0, -1):
                    if file_name[i] == ' ':
                        file_name = file_name[:i]
                        break

                for i in range(len(file_name) - 1, 0, -1):
                    if file_name[i] == '/':
                        file_name = file_name[i + 1:]
                        break
            else:
                file_name = "VSM Data"

            Final_png_path = self.Final_Hysteresis + '/' + '*.png'
            Final_png_list = glob.glob(Final_png_path)
            data_number_Final_png_list = len(Final_png_list)

            for i in range(0, data_number_Final_png_list, 1):
                for j in range(len(Final_png_list[i]) - 1, 0, -1):
                    if Final_png_list[i][j] == '/':
                        Final_png_list[i] = Final_png_list[i][j + 1:]
                        break

            def extract_number(f):
                try:
                    return int(f.split('K')[0])
                except:
                    return 0

            Final_png_list = sorted(Final_png_list, key=extract_number)

            # Final_png_RAW_path = self.Final_Hysteresis_RAW + '/' + '*.png'
            # Final_png_RAW_list = glob.glob(Final_png_RAW_path)
            # data_number_Final_png_RAW_list = len(Final_png_RAW_list)

            # for i in range(0, data_number_Final_png_RAW_list, 1):
            #     for j in range(len(Final_png_RAW_list[i]) - 1, 0, -1):
            #         if Final_png_RAW_list[i][j] == '/':
            #             Final_png_RAW_list[i] = Final_png_RAW_list[i][j + 1:]
            #             break

            # Final_png_RAW_list = sorted(Final_png_RAW_list, key=extract_number)

            inital_x = 0
            initial_y = 0.89
            image_width = 2.15
            x_offset = 2.24
            y_offset = 1.67
            iteration = 0

            blank_slide_layout = prs.slide_layouts[6]
            slide = prs.slides.add_slide(blank_slide_layout)
            text_frame = slide.shapes.add_textbox(Inches(5.67), Inches(0.06), Inches(2.03), Inches(0.57))
            text_frame = text_frame.text_frame
            p = text_frame.paragraphs[0]
            run = p.add_run()
            run.text = str(file_name)
            font = run.font
            font.name = 'Calibri'
            font.size = Pt(28)

            for k in range(0, data_number_Final_png_list, 1):
                VSM = self.Final_Hysteresis + '/' + Final_png_list[k]
                if iteration == 6:
                    inital_x = 0
                    initial_y = initial_y + y_offset
                    iteration = 0
                VSM_image = slide.shapes.add_picture(VSM, Inches(inital_x), Inches(initial_y), Inches(image_width))
                inital_x += x_offset
                iteration += 1

            prs.save(self.ProcCal + '/slope_removal.pptx')

            # blank_slide_layout = prs_RAW.slide_layouts[6]
            # slide = prs_RAW.slides.add_slide(blank_slide_layout)
            # text_frame = slide.shapes.add_textbox(Inches(5.67), Inches(0.06), Inches(2.03), Inches(0.57))
            # text_frame = text_frame.text_frame
            # p = text_frame.paragraphs[0]
            # run = p.add_run()
            # run.text = str(file_name)
            # font = run.font
            # font.name = 'Calibri'
            # font.size = Pt(28)
            #
            # inital_x = 0
            # initial_y = 0.89
            # iteration = 0

            # for k in range(0, data_number_Final_png_RAW_list, 1):
            #     VSM = self.Final_Hysteresis_RAW + '/' + Final_png_RAW_list[k]
            #     if iteration == 6:
            #         inital_x = 0
            #         initial_y = initial_y + y_offset
            #         iteration = 0
            #     VSM_image = slide.shapes.add_picture(VSM, Inches(inital_x), Inches(initial_y), Inches(image_width))
            #     inital_x += x_offset
            #     iteration += 1
            #
            # prs_RAW.save(self.ProcCal + '/raw_hysteresis.pptx')

        except Exception as e:
            print(f"Error generating PowerPoint: {e}")


# Main execution
if __name__ == "__main__":
    import sys
    from PyQt6.QtWidgets import QApplication

    app = QApplication(sys.argv)
    window = VSM_Data_Processing()
    window.show()
    sys.exit(app.exec())