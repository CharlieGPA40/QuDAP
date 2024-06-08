from PyQt6.QtWidgets import (
    QMessageBox, QButtonGroup, QWidget, QRadioButton, QGroupBox, QFileDialog, QVBoxLayout, QLabel, QHBoxLayout
, QDialog, QPushButton, QComboBox, QLineEdit, QScrollArea)
from PyQt6.QtGui import QIcon, QFont
from PyQt6.QtCore import pyqtSignal, Qt, QTimer, QThread
import sys
import pandas as pd
import matplotlib

matplotlib.use('QtAgg')
from matplotlib.backends.backend_qt5agg import FigureCanvasQTAgg as FigureCanvas
from matplotlib.backends.backend_qt5agg import NavigationToolbar2QT as NavigationToolbar
from matplotlib.figure import Figure
import matplotlib.patches as patches
import numpy as np
import traceback
import os


from pptx import Presentation
from pptx.util import Inches, Pt

class UserDefineFittingWindow(QDialog):
    def __init__(self):
        super().__init__()
        try:

            self.initUI()
        except Exception as e:
            QMessageBox.warning(self, "Error", str(e))
            return

    def initUI(self):
        self.setWindowTitle('User Defined Fitting Window')
        self.setGeometry(150, 150, 300, 200)

        self.layout = QVBoxLayout()

        self.label = QLabel('Enter some text:', self)
        self.layout.addWidget(self.label)

        self.textbox = QLineEdit(self)
        self.layout.addWidget(self.textbox)

        self.btn_close = QPushButton('Close', self)
        self.btn_close.clicked.connect(self.close)
        self.layout.addWidget(self.btn_close)

        self.setLayout(self.layout)



class MplCanvas(FigureCanvas):
    def __init__(self, parent=None, width=5, height=4, dpi=300, polar=False):
        self.fig = Figure(figsize=(width, height), dpi=dpi)
        # if polar:
        self.fig.clear()
        self.ax = self.fig.add_subplot(111, projection='polar' if polar else 'rectilinear')
        # else:
        #     self.fig.clear()
        #     self.ax = self.fig.add_subplot(111)
        super(MplCanvas, self).__init__(self.fig)


class General(QWidget):

    def __init__(self):
        super().__init__()

        try:
            self.autofittingWidget = False
            self.processbuttonWidget = False
            self.processWidget = False
            self.isInit = False
            self.plot_index = 0
            self.init_ui()

        except Exception as e:
            QMessageBox.warning(self, "Error", str(e))
            return


    def init_ui(self):
        try:
            if self.isInit == False:
                self.isInit = True
                titlefont = QFont("Arial", 20)
                self.font = QFont("Arial", 13)
                self.setStyleSheet("background-color: white;")
                # self.master_layout = QVBoxLayout()
                # # Create main vertical layout with centered alignment
                # self.scrollArea = QScrollArea(self)
                # self.scrollArea.setWidgetResizable(True)
                # self.scrollArea.setMinimumSize(800, 1000)

                # self.scrollContent = QWidget()
                # self.main_layout = QVBoxLayout(self.scrollContent)
                self.SHG_data_Processing_main_layout = QVBoxLayout(self)
                self.SHG_data_Processing_main_layout.setAlignment(Qt.AlignmentFlag.AlignCenter)

                #  ---------------------------- PART 1 --------------------------------
                self.SHG_General_label = QLabel("SHG Data Processing")
                self.SHG_General_label.setFont(titlefont)
                self.SHG_General_label.setStyleSheet("""
                                                    QLabel{
                                                        background-color: white;
                                                        }
                                                        """)
                #  ---------------------------- PART 2 --------------------------------
                self.file_selection_group_box = QGroupBox("Directory Selection")
                self.file_selection_section_layout = QHBoxLayout()
                self.file_selection_main_layout = QVBoxLayout()

                self.file_selection_label = QLabel('Please select the data directory:')
                self.file_selection_label.setFont(self.font)
                self.file_selection_button = QPushButton('Select')
                self.file_selection_button.setStyleSheet("""
                                   QPushButton {
                                       background-color: #3498DB; /* Green background */
                                       color: white; /* White text */
                                       border-style: solid;
                                       border-color: #3498DB;
                                       border-width: 2px;
                                       border-radius: 10px; /* Rounded corners */
                                       padding: 5px;
                                       min-height: 2px;
                                       min-width: 50px;
                                   }
                                   QPushButton:hover {
                                       background-color: #2874A6  ; /* Slightly darker green */
                                   }
                                   QPushButton:pressed {
                                       background-color: #85C1E9; /* Even darker green */
                                   }
                               """)
                self.file_selection_display_label = QLabel('Directory Selected: N/A')
                self.file_selection_display_label.setFont(QFont("Arial", 9))
                self.file_selection_button.clicked.connect(self.showDialog)
                self.file_selection_section_layout.addStretch(1)
                self.file_selection_section_layout.addWidget(self.file_selection_label, 2)
                self.file_selection_section_layout.addStretch(2)
                self.file_selection_section_layout.addWidget(self.file_selection_button, 1)
                self.file_selection_section_layout.addStretch(1)

                self.file_selection_main_layout.addLayout(self.file_selection_section_layout)
                self.file_selection_main_layout.addWidget(self.file_selection_display_label,
                                                          alignment=Qt.AlignmentFlag.AlignCenter)
                self.file_selection_main_layout.setAlignment(Qt.AlignmentFlag.AlignCenter)
                self.file_selection_group_box.setLayout(self.file_selection_main_layout)

                self.SHG_data_Processing_main_layout.addWidget(self.SHG_General_label, alignment=Qt.AlignmentFlag.AlignTop)
                self.SHG_data_Processing_main_layout.addWidget(self.file_selection_group_box)
                # self.scrollArea.setWidget(self.scrollContent)

                # Add the scroll area to the main layout
                # self.master_layout.addWidget(self.scrollArea)
                print('End')
        except Exception as e:
            QMessageBox.warning(self, "Error", str(e))
            return

    def showDialog(self):
        folder = QFileDialog.getExistingDirectory(self, "Select Folder")
        if folder:
            print(folder)
            folder = folder + "/"
            self.folder = folder
            self.file_selection_display_label.setText("Current Folder: {}".format(self.folder))
            self.auto_fitting()

    def auto_fitting(self):

        #  ---------------------------- PART 3&4 --------------------------------
        self.auto_fitting_selection_layout = QHBoxLayout()
        self.Auto_mode_group_box = QGroupBox("Auto Mode *")
        self.Fitting_mode_group_box = QGroupBox("Data Fitting *")
        #  ---------------------------- PART 3 --------------------------------
        self.auto_mode_layout = QVBoxLayout()
        self.auto_mode_yes_radio_buttom = QRadioButton("Auto")
        self.auto_mode_yes_radio_buttom.setFont(self.font)
        self.auto_mode_yes_radio_buttom.setChecked(True)
        self.auto_mode_no_radio_buttom = QRadioButton("Manual")
        self.auto_mode_no_radio_buttom.setFont(self.font)

        self.auto_mode_layout.addWidget(self.auto_mode_yes_radio_buttom)
        self.auto_mode_layout.addWidget(self.auto_mode_no_radio_buttom)

        self.buttonGroup = QButtonGroup()
        self.buttonGroup.addButton(self.auto_mode_yes_radio_buttom)
        self.buttonGroup.addButton(self.auto_mode_no_radio_buttom)

        self.auto_mode_yes_radio_buttom.toggled.connect(self.updateModeSelection)
        self.auto_mode_no_radio_buttom.toggled.connect(self.updateModeSelection)

        self.Auto_mode_group_box.setLayout(self.auto_mode_layout)
        self.auto_fitting_selection_layout.addWidget(self.Auto_mode_group_box)

        #  ---------------------------- PART 4 -------------------------------
        self.fitting_mode_layout = QVBoxLayout()
        self.fitting_mode_predef_radio_buttom = QRadioButton("Pre-defined")
        self.fitting_mode_predef_radio_buttom.setFont(self.font)
        self.fitting_mode_usr_radio_buttom = QRadioButton("User-defined")
        self.fitting_mode_usr_radio_buttom.setFont(self.font)
        self.fitting_mode_no_radio_buttom = QRadioButton("No")
        self.fitting_mode_no_radio_buttom.setFont(self.font)
        self.fitting_mode_no_radio_buttom.setChecked(True)

        self.fitting_mode_predef_layout = QHBoxLayout()
        self.fitting_mode_predef_combo = QComboBox()
        self.fitting_mode_predef_combo.setStyleSheet("""
                    QComboBox {
                        padding: 5px;
                        background-color: white;
                        border: 2px solid #c0c0c0;
                        border-radius: 4px;
                        }
                    QComboBox::item:selected {
                        background-color: #FFFFFF;  /* Background color for selected item */
                        color: #7CACEC
                    }

                    QComboBox::drop-down {
                        subcontrol-origin: padding;
                        subcontrol-position: top right;
                        width: 20px;   /* Width of the arrow button */
                        border-left-width: 1px;
                        border-left-color: darkgray;
                        border-left-style: solid; /* just a single line at the left */
                        border-top-right-radius: 3px; /* same radius as the QComboBox */
                        border-bottom-right-radius: 3px;
                    }
                    QComboBox::down-arrow {
                        image: url(GUI/Icon/chevron-down.svg); /* Set your own icon for the arrow */
                    }
                    QComboBox::down-arrow:on { /* When the combo box is open */
                        top: 1px;
                        left: 1px;
                    }
                """)
        self.fitting_mode_predef_combo.setFont(self.font)
        self.fitting_mode_predef_combo.addItems(["Please select a model"])
        # self.fitting_mode_predef_combo.setEnabled(True)
        self.fitting_mode_predef_button = QPushButton('Select')
        self.fitting_mode_predef_button.setStyleSheet("""
                                                                   QPushButton {
                                                                       background-color: #CAC9Cb; /* Green background */
                                                                       color: black; /* White text */
                                                                       border-style: solid;
                                                                       border-color: #CAC9Cb;
                                                                       border-width: 2px;
                                                                       border-radius: 10px; /* Rounded corners */
                                                                       padding: 5px;
                                                                       min-height: 1px;
                                                                       min-width: 50px;
                                                                   }
                                                                   QPushButton:hover {
                                                                       background-color: #5F6A6A; /* Slightly darker green */
                                                                   }
                                                                   QPushButton:pressed {
                                                                       background-color: #979A9A; /* Even darker green */
                                                                   }
                                                               """)
        self.fitting_mode_predef_button.setFont(self.font)

        # self.fitting_mode_predef_button.setEnabled(True)
        self.fitting_mode_predef_layout.addWidget(self.fitting_mode_predef_radio_buttom)
        self.fitting_mode_predef_layout.addWidget(self.fitting_mode_predef_combo)
        self.fitting_mode_predef_layout.addWidget(self.fitting_mode_predef_button)

        self.fitting_mode_usrdef_layout = QHBoxLayout()
        self.fitting_mode_usrdef_button = QPushButton('Define')
        self.fitting_mode_usrdef_button.setStyleSheet("""
                                                                   QPushButton {
                                                                       background-color: #CAC9Cb; /* Green background */
                                                                       color: black; /* White text */
                                                                       border-style: solid;
                                                                       border-color: #CAC9Cb;
                                                                       border-width: 2px;
                                                                       border-radius: 10px; /* Rounded corners */
                                                                       padding: 5px;
                                                                       min-height: 1px;
                                                                       min-width: 50px;
                                                                   }
                                                                   QPushButton:hover {
                                                                       background-color: #5F6A6A; /* Slightly darker green */
                                                                   }
                                                                   QPushButton:pressed {
                                                                       background-color: #979A9A; /* Even darker green */
                                                                   }
                                                               """)
        self.fitting_mode_usrdef_button.setFont(self.font)
        # self.fitting_mode_usrdef_button.setEnabled(True)
        self.fitting_mode_usrdef_layout.addWidget(self.fitting_mode_usr_radio_buttom)
        self.fitting_mode_usrdef_layout.addWidget(self.fitting_mode_usrdef_button)

        self.fitting_mode_layout.addLayout(self.fitting_mode_predef_layout)
        self.fitting_mode_layout.addLayout(self.fitting_mode_usrdef_layout)
        self.fitting_mode_layout.addWidget(self.fitting_mode_no_radio_buttom)

        self.Fitting_mode_group_box.setLayout(self.fitting_mode_layout)
        self.auto_fitting_selection_layout.addWidget(self.Fitting_mode_group_box)

        self.fitting_buttonGroup = QButtonGroup()
        self.fitting_buttonGroup.addButton(self.fitting_mode_predef_radio_buttom)
        self.fitting_buttonGroup.addButton(self.fitting_mode_usr_radio_buttom)
        self.fitting_buttonGroup.addButton(self.fitting_mode_no_radio_buttom)

        self.fitting_mode_predef_radio_buttom.toggled.connect(self.updateFitSelection)
        self.fitting_mode_usr_radio_buttom.toggled.connect(self.updateFitSelection)
        self.fitting_mode_no_radio_buttom.toggled.connect(self.updateFitSelection)

        #  ---------------------------- Main Layout --------------------------------
        # Add widgets to the main layout with centered alignment
        self.SHG_data_Processing_main_layout.addLayout(self.auto_fitting_selection_layout)
        # self.scrollArea.setWidget(self.scrollContent)

        # Add the scroll area to the main layout
        # self.master_layout.addWidget(self.scrollArea)
        self.process_layout = QHBoxLayout()
        self.rst_button = QPushButton('Reset')
        self.rst_button.setFont(self.font)
        self.rst_button.setStyleSheet("""
                                                                               QPushButton {
                                                                                   background-color: #CAC9Cb; /* Green background */
                                                                                   color: black; /* White text */
                                                                                   border-style: solid;
                                                                                   border-color: #CAC9Cb;
                                                                                   border-width: 2px;
                                                                                   border-radius: 10px; /* Rounded corners */
                                                                                   padding: 5px;
                                                                                   min-height: 1px;
                                                                                   min-width: 50px;
                                                                               }
                                                                               QPushButton:hover {
                                                                                   background-color: #5F6A6A; /* Slightly darker green */
                                                                               }
                                                                               QPushButton:pressed {
                                                                                   background-color: #979A9A; /* Even darker green */
                                                                               }
                                                                           """)
        self.rst_button.clicked.connect(self.rstpage)
        self.Process_button = QPushButton('Process')
        self.Process_button.setStyleSheet("""
                                   QPushButton {
                                       background-color: #3498DB; /* Green background */
                                       color: white; /* White text */
                                       border-style: solid;
                                       border-color: #3498DB;
                                       border-width: 2px;
                                       border-radius: 10px; /* Rounded corners */
                                       padding: 5px;
                                       min-height: 2px;
                                       min-width: 50px;
                                   }
                                   QPushButton:hover {
                                       background-color: #2874A6  ; /* Slightly darker green */
                                   }
                                   QPushButton:pressed {
                                       background-color: #85C1E9; /* Even darker green */
                                   }
                               """)
        self.Process_button.clicked.connect(self.process_data)
        self.Process_button.setFont(self.font)
        self.process_layout.addStretch(3)
        self.process_layout.addWidget(self.rst_button)
        self.process_layout.addWidget(self.Process_button)
        self.autofittingWidget = True
        self.processbuttonWidget = True
        self.SHG_data_Processing_main_layout.addLayout(self.process_layout)
        # self.scrollArea.setWidget(self.scrollContent)

        # Add the scroll area to the main layout
        # self.master_layout.addWidget(self.scrollArea)

        self.setLayout(self.SHG_data_Processing_main_layout)

        #  ---------------------------- Style Sheet --------------------------------

    def PresetModel(self):
        self.current_model = self.fitting_mode_predef_combo.currentText()
        print(self.current_model)

    def updateModeSelection(self):
        selectedModeButton = self.buttonGroup.checkedButton()
        if selectedModeButton:
            self.modeSelected = True
            self.auto = str(selectedModeButton.text())
        else:
            self.modeSelected = False
            self.auto = 'None'

    def updateFitSelection(self):
        selectedButton = self.fitting_buttonGroup.checkedButton()
        if selectedButton:
            self.fitSelected = True
            self.FitSeleciton = selectedButton.text()
            if self.FitSeleciton == 'Pre-defined':
                self.fitting_mode_predef_combo.setEnabled(True)
                self.fitting_mode_predef_button.setEnabled(True)
                self.fitting_mode_usrdef_button.setEnabled(False)
                self.fitting_mode_predef_button.clicked.connect(self.PresetModel)
            elif self.FitSeleciton == 'User-defined':
                self.fitting_mode_predef_combo.setEnabled(False)
                self.fitting_mode_predef_button.setEnabled(False)
                self.fitting_mode_usrdef_button.setEnabled(True)
                self.fitting_mode_usrdef_button.clicked.connect(self.showFittingPopup)
            else:
                self.fitting_mode_predef_combo.setEnabled(False)
                self.fitting_mode_predef_button.setEnabled(False)
                self.fitting_mode_usrdef_button.setEnabled(False)
        else:
            self.fitSelected = False

    def process_data(self):
        try:
            self.updateModeSelection()
            self.updateFitSelection()
            if self.modeSelected and self.fitSelected:
                #  ---------------------------- PART 5&6 --------------------------------
                self.data_processing_layout = QHBoxLayout()
                self.log_group_box = QGroupBox("Experimental Log")
                self.Fitting_mode_group_box = QGroupBox("Graph")
                #  ---------------------------- PART 3 --------------------------------
                self.Parametre = pd.read_csv(self.folder + "Experimental_Parameters.txt", header=None, sep=':', engine='c')
                Date = self.Parametre.iat[0, 1]
                self.date_label = QLabel('Date: ' + str(Date))
                self.date_label.setFont(self.font)
                file_name = self.Parametre.iat[1, 1]
                self.sample_label = QLabel('Sample: ' + str(file_name))
                self.sample_label.setFont(self.font)
                Measure_Type = self.Parametre.iat[2, 1]
                self.measure_type_label = QLabel('Measurement Type: ' + str(Measure_Type))
                self.measure_type_label.setFont(self.font)
                Light_angle = self.Parametre.iat[3, 1]
                self.light_angle_label = QLabel('Incident Angle: ' + str(Light_angle))
                self.light_angle_label.setFont(self.font)
                power = self.Parametre.iat[4, 1]
                self.power_label = QLabel('Power (mW): ' + str(power))
                self.power_label.setFont(self.font)
                start_angle = self.Parametre.iat[5, 1]
                self.start_angle_label = QLabel('Start Angle (degree): ' + str(start_angle))
                self.start_angle_label.setFont(self.font)
                end_angle = self.Parametre.iat[6, 1]
                self.end_angle_label = QLabel('Termination Angle (degree): ' + str(end_angle))
                self.end_angle_label.setFont(self.font)
                step_size = self.Parametre.iat[7, 1]
                self.step_angle_label = QLabel('Angle Step Size (degree): ' + str(step_size))
                self.step_angle_label.setFont(self.font)
                self.step_size = int(step_size)
                polarization = self.Parametre.iat[8, 1]
                self.polarization = polarization
                self.polarization_label = QLabel('Polarization Configuration: ' + str(polarization))
                self.polarization_label.setFont(self.font)
                exp_time = self.Parametre.iat[9, 1]
                self.exp_time_label = QLabel('Exposure Time (s): ' + str(exp_time))
                self.exp_time_label.setFont(self.font)
                EMGain = self.Parametre.iat[10, 1]
                self.EMGain_label = QLabel('EMGain: ' + str(EMGain))
                self.EMGain_label.setFont(self.font)
                Accumulation = self.Parametre.iat[11, 1]
                self.acc_label = QLabel('Accumulation: ' + str(Accumulation))
                self.acc_label.setFont(self.font)
                Start_temp = self.Parametre.iat[12, 1]
                self.start_temp_label = QLabel('Initial Temperature (K): ' + str(Start_temp))
                self.start_temp_label.setFont(self.font)
                End_temp = self.Parametre.iat[13, 1]
                self.end_temp_label = QLabel('Termination Temperature (K): ' + str(End_temp))
                self.end_temp_label.setFont(self.font)
                Step_temp = self.Parametre.iat[14, 1]
                self.step_temp_label = QLabel('Step Temperature (K): ' + str(Step_temp))
                self.step_temp_label.setFont(self.font)

                # If need to remove the space
                file_name = str(self.Parametre.iat[1, 1]).replace(" ", "")
                self.file_name = file_name
                self.log_layout = QVBoxLayout()
                self.log_layout.addWidget(self.date_label)
                self.log_layout.addWidget(self.sample_label)
                self.log_layout.addWidget(self.measure_type_label)
                self.log_layout.addWidget(self.light_angle_label)
                self.log_layout.addWidget(self.power_label)
                self.log_layout.addWidget(self.start_angle_label)
                self.log_layout.addWidget(self.end_angle_label)
                self.log_layout.addWidget(self.step_angle_label)
                self.log_layout.addWidget(self.polarization_label)
                self.log_layout.addWidget(self.exp_time_label)
                self.log_layout.addWidget(self.EMGain_label)
                self.log_layout.addWidget(self.acc_label)
                self.log_layout.addWidget(self.start_temp_label)
                self.log_layout.addWidget(self.end_temp_label)
                self.log_layout.addWidget(self.step_temp_label)

                self.log_group_box.setLayout(self.log_layout)
                self.data_processing_layout.addWidget(self.log_group_box)
                figure_Layout = QVBoxLayout()
                self.canvas = MplCanvas(self, width=5, height=4, dpi=100, polar=False)
                toolbar = NavigationToolbar(self.canvas, self)
                toolbar.setStyleSheet("""
                                                               QWidget {
                                                                   border: None;
                                                               }
                                                           """)

                plot_entry_layout = QHBoxLayout()
                self.center_x_y_label = QLabel("Center (x, y): ")
                self.center_x_y_label.setFont(self.font)
                self.center_front_p_label = QLabel("(")
                self.center_front_p_label.setFont(self.font)
                self.center_x_entry_box = QLineEdit()
                self.center_comma_label = QLabel(",")
                self.center_comma_label.setFont(self.font)
                self.center_y_entry_box = QLineEdit()
                self.center_end_p_label = QLabel(")")
                self.center_end_p_label.setFont(self.font)

                self.box_size_label = QLabel("Box Size (pixels): ")
                self.box_size_label.setFont(self.font)
                self.box_size_entry = QLineEdit()

                plot_entry_layout.addWidget(self.center_x_y_label)
                plot_entry_layout.addWidget(self.center_front_p_label)
                plot_entry_layout.addWidget(self.center_x_entry_box)
                plot_entry_layout.addWidget(self.center_comma_label)
                plot_entry_layout.addWidget(self.center_y_entry_box)
                plot_entry_layout.addWidget(self.center_end_p_label)
                plot_entry_layout.addStretch(1)
                plot_entry_layout.addWidget(self.box_size_label)
                plot_entry_layout.addWidget(self.box_size_entry)

                plot_button_layout = QHBoxLayout()
                self.prev_button = QPushButton("Previous")
                self.prev_button.clicked.connect(self.show_previous_plot)
                self.prev_button.setStyleSheet("""
                                                   QPushButton {
                                                       background-color: #CAC9Cb; /* Green background */
                                                       color: black; /* White text */
                                                       border-style: solid;
                                                       border-color: #CAC9Cb;
                                                       border-width: 2px;
                                                       border-radius: 10px; /* Rounded corners */
                                                       padding: 5px;
                                                       min-height: 1px;
                                                       min-width: 10px;
                                                   }
                                                   QPushButton:hover {
                                                       background-color: #5F6A6A; /* Slightly darker green */
                                                   }
                                                   QPushButton:pressed {
                                                       background-color: #979A9A; /* Even darker green */
                                                   }
                                               """)
                self.next_button = QPushButton("Next")
                self.next_button.setStyleSheet("""
                                                   QPushButton {
                                                       background-color: #CAC9Cb; /* Green background */
                                                       color: black; /* White text */
                                                       border-style: solid;
                                                       border-color: #CAC9Cb;
                                                       border-width: 2px;
                                                       border-radius: 10px; /* Rounded corners */
                                                       padding: 5px;
                                                       min-height: 1px;
                                                       min-width: 10px;
                                                   }
                                                   QPushButton:hover {
                                                       background-color: #5F6A6A; /* Slightly darker green */
                                                   }
                                                   QPushButton:pressed {
                                                       background-color: #979A9A; /* Even darker green */
                                                   }
                                               """)
                self.next_button.clicked.connect(self.show_next_plot)

                plot_button_layout.addWidget(self.prev_button)
                plot_button_layout.addWidget(self.next_button)

                figure_Layout.addWidget(toolbar, alignment=Qt.AlignmentFlag.AlignCenter)
                figure_Layout.addWidget(self.canvas, alignment=Qt.AlignmentFlag.AlignCenter)
                figure_Layout.addLayout(plot_entry_layout)
                figure_Layout.addLayout(plot_button_layout)

                self.Fitting_mode_group_box.setLayout(figure_Layout)
                self.data_processing_layout.addWidget(self.Fitting_mode_group_box)
                self.processWidget = True
                self.SHG_data_Processing_main_layout.addLayout(self.data_processing_layout)

                self.setLayout(self.SHG_data_Processing_main_layout)
                self.degree = 130
                self.SHG_Raw = np.loadtxt(self.folder + file_name + "_{}deg".format(self.degree) + ".txt", dtype=int,
                                          delimiter=',')
                self.title = str(file_name) + '' + str(Measure_Type) + '' + str(Light_angle) + '\n' + str(power) + ('mW Exposure Time ') + str(exp_time) + 's Averaging ' + str(Accumulation)
                if self.auto == 'Auto':
                    max_sum = -np.inf
                    max_region = None
                    best_size = 0
                    min_region_size = 10
                    max_region_size = 150
                    abs_SHG_Raw = np.abs(self.SHG_Raw)
                    background_pixel = np.sum(self.SHG_Raw[450:512, 450:512]) / (512 - 450) ** 2

                    for i in range(240, 280, 1):
                        for j in range(240, 280, 1):
                            center_region = abs_SHG_Raw[i:i + min_region_size, j:j + min_region_size]
                            pixel_region_sum = np.sum(center_region)
                            pixel_region_sum = pixel_region_sum - background_pixel * (min_region_size ** 2)
                            if pixel_region_sum > max_sum:
                                max_sum = pixel_region_sum
                                best_size = min_region_size
                                max_region = (i, j, best_size, max_sum)

                    start_i, start_j, region_size_first, region_sum = max_region
                    center_i = start_i + region_size_first // 2
                    center_j = start_j + region_size_first // 2

                    max_center_i_start = int(center_i - max_region_size / 2)
                    max_center_i_end = int(center_i + max_region_size / 2)
                    max_center_j_start = int(center_j - max_region_size / 2)
                    max_center_j_end = int(center_j + max_region_size / 2)

                    min_center_i_start = int(center_i - min_region_size / 2)
                    min_center_i_end = int(center_i + min_region_size / 2)
                    min_center_j_start = int(center_j - min_region_size / 2)
                    min_center_j_end = int(center_j + min_region_size / 2)

                    max_sum_in_region = abs_SHG_Raw[max_center_i_start:max_center_i_end,
                                        max_center_j_start:max_center_j_end]
                    min_sum_in_region = abs_SHG_Raw[min_center_i_start:min_center_i_end,
                                        min_center_j_start:min_center_j_end]

                    max_pixel_region_sum = np.sum(max_sum_in_region) - background_pixel * (max_region_size ** 2)
                    min_pixel_region_sum = np.sum(min_sum_in_region) - background_pixel * (min_region_size ** 2)
                    difference_max_min = max_pixel_region_sum - min_pixel_region_sum
                    pixel_region_old = min_pixel_region_sum
                    for region_size in range(min_region_size, max_region_size, 2):
                        center_i_start = int(center_i - region_size / 2)
                        center_i_end = int(center_i + region_size / 2)
                        center_j_start = int(center_j - region_size / 2)
                        center_j_end = int(center_j + region_size / 2)

                        region = abs_SHG_Raw[center_i_start:center_i_end, center_j_start:center_j_end]
                        pixel_region_sum = np.sum(region) - background_pixel * (region_size ** 2)

                        if pixel_region_sum - pixel_region_old > 0.2 * difference_max_min:
                            max_sum = pixel_region_sum
                            pixel_region_old = pixel_region_sum
                            best_size = region_size
                            max_region = (center_i, center_j, best_size, max_sum)

                    center_i, center_j, region_size, region_sum = max_region
                    start_i = center_i - region_size // 2
                    start_j = center_j - region_size // 2
                    self.start_i = start_i
                    self.start_j = start_j
                    self.region_size = region_size

                    self.center_x_entry_box.setEnabled(False)
                    self.center_y_entry_box.setEnabled(False)
                    self.box_size_entry.setEnabled(False)
                    self.center_x_entry_box.setText(str(center_i))
                    self.center_y_entry_box.setText(str(center_j))
                    self.box_size_entry.setText(str(region_size))
                    self.prev_button.setEnabled(False)

                    # self.canvas.figure.clf()
                    # self.canvas = MplCanvas(self, width=5, height=4, dpi=100, polar=False)
                    self.canvas.ax.imshow(self.SHG_Raw, aspect='auto', vmin=0, vmax=5000)
                    # self.canvas.figure.colorbar(label='{} Polarization'.format(polarization))
                    self.canvas.ax.scatter(center_i, center_j, s=30, color='tomato',
                                           marker='x')
                    self.canvas.ax.set_title(self.title + ' at {} Degree'.format(self.degree), pad=10, wrap=True)
                    self.canvas.figure.gca().add_patch(patches.Rectangle((start_i, start_j), region_size, region_size,
                                                                         edgecolor='white', facecolor='none', linewidth=1))
                    self.canvas.figure.savefig(self.folder + "Preview_Figure_at_130_Deg.png")
                    self.canvas.draw()
                    self.next_button.setText("Polar Plot")
                else:
                    self.box_size_entry.setPlaceholderText("Enter Box Size")
                    self.canvas = MplCanvas(self, width=5, height=4, dpi=100, polar=False)
                    self.canvas.ax.imshow(self.SHG_Raw, aspect='auto', vmin=0, vmax=5000)
                    self.canvas.ax.set_title(self.title + ' at {} Degree'.format(self.degree), pad=10, wrap=True)
                    def onclick(event):
                        if event.dblclick:
                            global cur_x, cur_y, click
                            cur_x = event.xdata
                            cur_y = event.ydata
                            click = event.button
                            if click == 1:
                                print('Closing')

                                self.center_x_entry_box.setEnabled(False)
                                self.center_y_entry_box.setEnabled(False)
                                self.center_x = int(cur_x)
                                self.center_y = int(cur_y)
                                self.center_x_entry_box.setText(str(int(cur_x)))
                                self.center_y_entry_box.setText(str(int(cur_y)))

                    connection_id = self.canvas.figure.canvas.mpl_connect('button_press_event', onclick)
                    self.canvas.figure.savefig(self.folder + "Preview_Figure_at_130_Deg.png")
                    self.canvas.draw()
                    self.prev_button.setEnabled(True)
                    self.next_button.setText("Submit")

            else:
                QMessageBox.warning(self, "Please try again!", "Select all the required option")
        except Exception as e:
            QMessageBox.warning(self, "Error", str(e))
            return



    def show_previous_plot(self):
        try:
            if self.plot_index == 0:
                self.prev_button.setEnabled(False)
            elif self.plot_index == 1:
                if self.auto == 'Manual':
                    self.box_size_entry.setPlaceholderText("Enter Box Size")
                    self.canvas = MplCanvas(self, width=5, height=4, dpi=100, polar=False)
                    self.canvas.ax.imshow(self.SHG_Raw, aspect='auto', vmin=0, vmax=5000)
                    self.canvas.ax.set_title(self.title + ' at {} Degree'.format(self.degree), pad=10, wrap=True)

                    def onclick(event):
                        if event.dblclick:
                            global cur_x, cur_y, click
                            cur_x = event.xdata
                            cur_y = event.ydata
                            click = event.button
                            if click == 1:
                                self.center_x_entry_box.setEnabled(False)
                                self.center_y_entry_box.setEnabled(False)
                                self.center_x = int(cur_x)
                                self.center_y = int(cur_y)
                                self.center_x_entry_box.setText(str(int(cur_x)))
                                self.center_y_entry_box.setText(str(int(cur_y)))

                    connection_id = self.canvas.figure.canvas.mpl_connect('button_press_event', onclick)
                    self.canvas.figure.savefig(self.folder + "Preview_Figure_at_130_Deg.png")
                    self.canvas.draw()
                    self.prev_button.setEnabled(True)
                    self.plot_index -= 1


        except Exception as e:
            QMessageBox.warning(self, "Error", str(e))

    def show_next_plot(self):
        try:
            if self.plot_index == 0:
                if self.auto == 'Auto':
                    self.polar_plot_extraction()
                else:
                    self.canvas.ax.clear()
                    self.canvas = MplCanvas(self, width=5, height=4, dpi=100, polar=False)
                    self.canvas.ax.imshow(self.SHG_Raw, aspect='auto', vmin=0, vmax=5000)
                    self.canvas.ax.set_title(self.title + ' at {} Degree'.format(self.degree), pad=10, wrap=True)

                    region_size = int(self.box_size_entry.displayText())
                    self.region_size = region_size
                    self.half_region_size = (np.ceil(region_size / 2)).astype(int)

                    self.canvas.ax.scatter(self.center_x, self.center_y, s=30, color='tomato',
                                           marker='x')

                    self.canvas.figure.gca().add_patch(
                        patches.Rectangle((self.center_x - self.half_region_size, self.center_y - self.half_region_size), region_size,
                                          region_size,
                                          edgecolor='white', facecolor='none', linewidth=1))
                    self.canvas.figure.savefig(self.folder + "Preview_Figure_at_130_Deg.png")
                    self.canvas.draw()
                    self.next_button.setText("Polar Plot")
                self.plot_index += 1
            elif self.plot_index == 1:
                if self.auto == 'Auto':
                    self.polar_plot_linear()
                else:
                    self.polar_plot_extraction()
                self.plot_index += 1
            elif self.plot_index == 2:
                if self.auto == 'Auto':
                    self.polar_plot_linear_correction()
                else:
                    self.polar_plot_linear()
                self.plot_index += 1
            elif self.plot_index == 3:
                if self.auto == 'Auto':
                    self.polar_plot_linear_neg_correction()
                    if self.FitSeleciton == 'Pre-defined' or self.FitSeleciton == 'User-defined':
                        self.plot_index = 5
                        self.next_button.setText("Fit")
                    else:
                        self.plot_index = 6
                        self.next_button.setText("Exp. PPT.")
                else:
                    self.polar_plot_linear_correction()
                    self.plot_index += 1
            elif self.plot_index == 4:
                self.polar_plot_linear_neg_correction()
                if self.FitSeleciton == 'Pre-defined' or self.FitSeleciton == 'User-defined':
                    self.plot_index = 5
                    self.next_button.setText("Fit")
                else:
                    self.plot_index = 6
                    self.next_button.setText("Exp. PPT.")
            elif self.plot_index == 5:
                self.Fitting()
                self.plot_index += 1

            elif self.plot_index == 6:
                self.export_PPT()
                self.plot_index = 0

        except Exception as e:
            QMessageBox.warning(self, "Error", str(e))

    def polar_plot_extraction(self):
        try:
            self.deg_file = []
            self.sig_file = []

            for degree in range(0, 365, self.step_size):
                self.deg_file = np.append(self.deg_file, degree)
                SHG_Raw = np.loadtxt(self.folder + self.file_name + "_{}deg".format(degree) + ".txt", dtype=int,
                                     delimiter=',')
                if self.auto == 'Manual':
                    region = SHG_Raw[(self.center_x - self.half_region_size): (self.center_x + self.half_region_size),
                             (self.center_y - self.half_region_size): (self.center_y + self.half_region_size)]
                else:
                    region = SHG_Raw[self.start_i: self.start_i + self.region_size, self.start_j: self.start_j + self.region_size]

                small_sum = sum(map(sum, region))
                large_sum = sum(map(sum, SHG_Raw))
                bkg_pixel = (large_sum - small_sum) / (512 ** 2 - self.region_size ** 2)
                sig = small_sum - bkg_pixel * self.region_size ** 2
                self.sig_file = np.append(self.sig_file, sig)

            self.sig_file = self.sig_file.astype(np.float64)
            max_lim = max(self.sig_file)
            min_lim = min(self.sig_file)
            self.deg_file = self.deg_file * np.pi / 180
            self.deg_file = self.deg_file.astype(np.float64)
            self.canvas.ax.clear()
            self.canvas = MplCanvas(self, width=5, height=4, dpi=100, polar=True)
            self.canvas.ax.plot(self.deg_file, self.sig_file, color='red')
            self.canvas.ax.scatter(self.deg_file, self.sig_file, color='orange')
            self.canvas.ax.set_ylim(bottom=min_lim, top=max_lim)
            self.canvas.figure.savefig(self.folder + "Raw_Polar_Plot.png")
            self.canvas.draw()

            csv_file_name = 'Processed_First_Data.csv'
            comb = pd.DataFrame(list(zip(self.deg_file, self.sig_file)))
            rec_data = pd.DataFrame()
            rec_data = pd.concat([rec_data, comb], ignore_index=True, axis=1)
            rec_data.to_csv(self.folder + csv_file_name, mode='a', index=False, encoding='utf-8-sig', header=None)
            self.next_button.setText("Linear Profile")
        except Exception as e:
            QMessageBox.warning(self, "Error", str(e))

    def polar_plot_linear(self):
        try:
            self.canvas.ax.clear()
            self.canvas = MplCanvas(self, width=5, height=4, dpi=100, polar=False)
            self.canvas.ax.plot(self.deg_file, self.sig_file, color='red')
            self.canvas.figure.savefig(self.folder + "Raw_Polar_Plot_Linear.png")
            self.canvas.draw()
            self.next_button.setText("Slope Correction")
        except Exception as e:
            QMessageBox.warning(self, "Error", str(e))

    def polar_plot_linear_correction(self):
        try:
            slope = (self.sig_file[-1] - self.sig_file[0]) / (self.deg_file[-1] - self.deg_file[0])
            const = self.sig_file[-1] - slope * self.deg_file[-1]
            # const = self.sig_file[-1] + slope * self.deg_file[-1]
            self.sig_file = self.sig_file - (slope * self.deg_file + const)
            # # sig_file = (sig_file / 30) / 380000
            csv_file_name = 'Linear_Slope_Correction.csv'
            comb = pd.DataFrame(list(zip(self.deg_file, self.sig_file)))
            rec_data = pd.DataFrame()
            rec_data = pd.concat([rec_data, comb], ignore_index=True, axis=1)
            rec_data.to_csv(self.folder + csv_file_name, mode='a', index=False, encoding='utf-8-sig', header=None)
            self.canvas.ax.clear()
            self.canvas = MplCanvas(self, width=5, height=4, dpi=100, polar=False)
            self.canvas.ax.plot(self.deg_file, self.sig_file, color='red')
            self.canvas.ax.scatter(self.deg_file, self.sig_file, color='orange')
            self.canvas.figure.savefig(self.folder + "Slope_removal.png")
            self.canvas.draw()
            self.next_button.setText("Negative Correction")
        except Exception as e:
            QMessageBox.warning(self, "Error", str(e))

    def polar_plot_linear_neg_correction(self):
            min_sig = min(self.sig_file)
            self.sig_file = self.sig_file - min_sig
            csv_file_name = 'Shift_Negative_To_Positive.csv'
            comb = pd.DataFrame(list(zip(self.deg_file, self.sig_file)))
            rec_data = pd.DataFrame()
            rec_data = pd.concat([rec_data, comb], ignore_index=True, axis=1)
            rec_data.to_csv(self.folder + csv_file_name, mode='w', index=False, encoding='utf-8-sig', header=None)
            max_lim = max(self.sig_file)
            min_lim = min(self.sig_file)
            #
            # fig, ax = pyplot.subplots(subplot_kw={'projection': 'polar'})
            # ax.plot(deg_file, sig_file, color='tomato')
            # ax.scatter(deg_file, sig_file, color='tomato')
            # ax.set_ylim(bottom=min_lim, top=max_lim)
            #
            # pyplot.tight_layout()
            # pyplot.savefig(folder_selected + "Figure_2.png")
            # pyplot.show()
            self.canvas.ax.clear()
            self.canvas = MplCanvas(self, width=5, height=4, dpi=100, polar=True)
            self.canvas.ax.plot(self.deg_file, self.sig_file, color='red')
            self.canvas.ax.set_ylim(bottom=min_lim, top=max_lim)
            self.canvas.ax.set_title(self.title + '{} Polarization'.format(self.polarization), pad=10, wrap=True)
            self.canvas.figure.savefig(self.folder + "Final_Processed_Data.png")
            self.canvas.draw()

    def Fitting(self):
           print('adding')
           self.next_button.setText("Exp. PPT.")
            #
            # #
            # fit = True
            # if fit == True:
            #     def shg_sin(params, x, data=None):
            #         A = params['A']
            #         a2 = params['a2']
            #         B = params['B']
            #         x0 = params['x0']
            #         model = (A * np.sin(a2 - 3 * (x - x0)) + B * np.sin(a2 + (x - x0))) ** 2
            #         if data is None:
            #             return model
            #         return model - data
            #
            #     def shg_cos(params, x, data=None):
            #
            #         A = params['A']
            #         a2 = params['a2']
            #         B = params['B']
            #         model = A * np.cos(a2 - 3 * x) + B * np.cos(a2 + 3 * x) * 2
            #         if data is None:
            #             return model
            #         return model - data
            #
            #     # Create a Parameters object
            #     params = lmfit.Parameters()
            #     params.add('A', value=-0.1)
            #     params.add('a2', value=-0.1)
            #     params.add('B', value=11)
            #     params.add('x0', value=0.1)
            #     result_sin = lmfit.minimize(shg_sin, params, args=(deg_file,), kws={'data': sig_file})
            #     sin_A = result_sin.params['A'].value
            #     sin_a2 = result_sin.params['a2'].value
            #     sin_B = result_sin.params['B'].value
            #     sin_x0 = result_sin.params['x0'].value
            #
            #     sin_A_err = result_sin.params['A'].stderr
            #     sin_a2_err = result_sin.params['a2'].stderr
            #     sin_B_err = result_sin.params['B'].stderr
            #     sin_x0_err = result_sin.params['x0'].stderr

            #     fig, ax = pyplot.subplots(subplot_kw={'projection': 'polar'})
            #
            #     ax.plot(deg_file, result_sin.residual + sig_file, color='#E74C3C', linewidth=2)
            #     ax.scatter(deg_file, sig_file, color='black', s=4)
            #     ax.set_ylim(bottom=min_lim * 1.1, top=max_lim * 1.1)
            #     pyplot.title(title + '{} Polarization'.format(polarization), pad=10, wrap=True)
            #     pyplot.tight_layout()
            #     pyplot.savefig(folder_selected + "Fitted_Data.png")
            #     pyplot.show()
            #
            #     df = pd.DataFrame()
            #     df_comb = pd.DataFrame(
            #         list(zip([sin_A], [sin_A_err], [sin_a2], [sin_a2_err], [sin_B], [sin_B_err], [sin_x0], [sin_x0_err])))
            #     df = pd.concat([df, df_comb], ignore_index=True, axis=1)
            #     df.to_csv(folder_selected + 'Fitted_Data.csv', index=False)
            #
            #     # df = pd.DataFrame()
            #     # df_comb = pd.DataFrame(list(zip([sin_A_err], [sin_a2_err], [sin_B_err], [sin_x0_err])))
            #     # df = pd.concat([df, df_comb], ignore_index=True, axis=1)
            #     # df.to_csv(folder_selected + 'Fitted_Data.csv', index=False)
    def export_PPT(self):
        try:
            if os.path.exists(self.folder + 'Results.pptx'):
                prs = Presentation(self.folder + 'Results.pptx')
                prs.slide_width = Inches(13.33)
                prs.slide_height = Inches(7.5)
            else:
                prs = Presentation()
                prs.slide_width = Inches(13.33)
                prs.slide_height = Inches(7.5)
                prs.save(self.folder + 'Results.pptx')
                prs = Presentation(self.folder + 'Results.pptx')
                prs.slide_width = Inches(13.33)
                prs.slide_height = Inches(7.5)

            SHG_Image = self.folder + "Preview_Figure_at_130_Deg.png"
            SHG_Signal = self.folder + "Final_Processed_Data.png"
            blank_slide_layout = prs.slide_layouts[6]
            slide = prs.slides.add_slide(blank_slide_layout)
            SHG_Image_img = slide.shapes.add_picture(SHG_Image, Inches(0.32), Inches(1.42), Inches(6.39))
            SHG_Signal_img = slide.shapes.add_picture(SHG_Signal, Inches(6.49), Inches(1.42), Inches(6.39))
            text_frame = slide.shapes.add_textbox(Inches(0.18), Inches(0.2), Inches(6.67), Inches(0.4))
            Data_frame = slide.shapes.add_textbox(Inches(11.19), Inches(0.2), Inches(2.04), Inches(0.4))
            text_frame = text_frame.text_frame
            Data_frame = Data_frame.text_frame
            p = text_frame.paragraphs[0]
            d = Data_frame.paragraphs[0]
            run = p.add_run()
            for i in range(len(self.folder) - 1, 0, -1):
                if self.folder[i] == '/':
                    folder_name_pptx = self.folder[i + 1:]
                    break
            run.text = str(folder_name_pptx)
            font = run.font
            font.name = 'Calibri'
            font.size = Pt(18)

            run_d = d.add_run()
            Date = self.Parameter.iat[0, 1]
            run_d.text = str(Date)
            font_d = run_d.font
            font_d.name = 'Calibri'
            font_d.size = Pt(18)
            prs.save(self.folder + 'Results.pptx')
        except Exception as e:
            QMessageBox.warning(self, "Error", str(e))

    def rstpage(self):
        try:
            self.folder = 'N/A'
            self.file_selection_display_label.setText("Current Folder: {}".format(self.folder))
            if self.autofittingWidget == True:
                self.autofittingWidget = False
                self.clearLayout(self.auto_fitting_selection_layout)
                self.SHG_data_Processing_main_layout.removeItem(self.auto_fitting_selection_layout)
            if self.processbuttonWidget == True:
                self.processbuttonWidget = False
                self.clearLayout(self.process_layout)
                self.SHG_data_Processing_main_layout.removeItem(self.process_layout)
            if self.processWidget == True:
                self.processWidget = False
                self.clearLayout(self.data_processing_layout)
                self.SHG_data_Processing_main_layout.removeItem(self.data_processing_layout)




        except Exception as e:
            QMessageBox.warning(self, "Error", str(e))
            return

    def clear_layout(self, layout):
        if layout is not None:
            print(f"Clearing layout with {layout.count()} items.")
            while layout.count():
                child = layout.takeAt(0)
                if child.widget() is not None:
                    print(f"Deleting widget: {child.widget()}")
                    child.widget().deleteLater()
                print(f"Precheck layout layout... {child.widget()}")
                if child.layout() is not None:
                    print(f"Clearing nested layout... {child.widget()}")
                    self.clear_layout(child.layout())
        print('cleared')
    def showFittingPopup(self):
        self.popup = UserDefineFittingWindow()
        self.popup.exec()

    def clearLayout(self, layout):
        if layout is not None:
            while layout.count():
                item = layout.takeAt(0)
                widget = item.widget()
                if widget is not None:
                    widget.deleteLater()
                elif item.layout() is not None:
                    self.clearLayout(item.layout())
