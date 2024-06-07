# -----------------------------------------------------------#
# Name: Chunli Tang                                          #
# School: Auburn University                                  #
# -----------------------------------------------------------#

import os
from tkinter import *
from matplotlib import pyplot
import matplotlib.patches as patches
import numpy as np
import pandas as pd
from tkinter import filedialog
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.shapes.autoshape import Shape as shape
from lmfit import Model
import lmfit
# from skimage import filters, measure, morphology, restoration


class SHG_Processing():
    def __init__(self):
        super().__init__()
        self.run()

    def run(self):
        root = Tk()
        folder_selected = filedialog.askdirectory(initialdir="/Users/labaccess/Library/CloudStorage/Box-Box/Jin Lab Shared Folder/SHG-RA data")
        # dir_list = os.listdir(folder_selected)
        root.withdraw()
        for i in range(len(folder_selected) - 1, 0, -1):
            if folder_selected[i] == '/':
                folder_name_pptx = folder_selected[i+1:]
                break

        folder_selected = folder_selected + "/"
        Parameter = pd.read_csv(folder_selected + "Experimental_Parameters.txt", header=None, sep=':', engine='c')
        print(Parameter)
        Date = Parameter.iat[0, 1]
        step_size = Parameter.iat[7, 1]
        step_size = int(step_size)
        file_name = str(Parameter.iat[1, 1]).replace(" ", "")

        avg_x = 0
        avg_y = 0
        iteration = 0

        deg_file_org = []
        sig_file_org = []

        # for degree in range(0, 10, step_size):
        degree = 130
        temp_temp = 80
        SHG_Raw = np.loadtxt(folder_selected + file_name + "_{}deg".format(degree) + ".txt", dtype=int, delimiter=',')
        auto = False
        if auto == True:
            region_size = 120
            # Variables to store the maximum intensity and its corresponding position
            max_intensity = -np.inf
            max_intensity_position = None

            # Iterate over all possible square regions
            for i in range(SHG_Raw.shape[0] - region_size + 1):
                for j in range(SHG_Raw.shape[1] - region_size + 1):
                    # Extract the current square region
                    region = SHG_Raw[i:i + region_size, j:j + region_size]

                    # Calculate the intensity of the region
                    intensity = np.sum(region)

                    # Check if the current intensity is the maximum
                    if intensity > max_intensity:
                        max_intensity = intensity
                        max_intensity_position = (i, j)

            # Get the coordinates of the square region with the maximum intensity
            top_left_row, top_left_col = max_intensity_position
            bottom_right_row = top_left_row + region_size - 1
            bottom_right_col = top_left_col + region_size - 1

            # Plot the heatmap with the square region of maximum intensity highlighted
            pyplot.imshow(SHG_Raw, aspect='auto', vmin=0,vmax=5000)
            polarization = Parameter.iat[8, 1]
            pyplot.colorbar(label='{} Polarization'.format(polarization))
            pyplot.scatter(top_left_col + region_size/2, top_left_row + region_size/2, s=30, color='tomato', marker='x')
            exposure_time = str(float(Parameter.iat[9, 1]))
            title = str(Parameter.iat[1, 1]) + '' + str(Parameter.iat[2, 1]) + '' + str(Parameter.iat[3, 1]) \
                    + '\n' + str(Parameter.iat[4, 1]) + 'mW Exposure Time ' + exposure_time + 's Averaging ' \
                    + str(int(Parameter.iat[11, 1]))
            pyplot.title(title + ' at {} Degree'.format(degree), pad=10, wrap=True)
            pyplot.gca().add_patch(pyplot.Rectangle((top_left_col, top_left_row), region_size, region_size,
                                              edgecolor='white', facecolor='none', linewidth=2))
            pyplot.savefig(folder_selected + "Figure_1.png")
            pyplot.show()

        else:
            fig, ax = pyplot.subplots()
            pyplot.imshow(SHG_Raw, vmin=0, vmax=5000)
            polarization = Parameter.iat[8, 1]
            pyplot.colorbar(label='{} Polarization'.format(polarization))
            exposure_time = str(float(Parameter.iat[9, 1]))
            title = str(Parameter.iat[1, 1]) + '' + str(Parameter.iat[2, 1]) + '' + str(Parameter.iat[3, 1]) \
                    + '\n' + str(Parameter.iat[4, 1]) + 'mW Exposure Time ' + exposure_time + 's Averaging ' \
                    + str(int(Parameter.iat[11, 1]))
            pyplot.title(title + ' at {} Degree'.format(degree), pad=10, wrap=True)
            def onclick(event):
                if event.dblclick:
                    global cur_x, cur_y, click
                    cur_x = event.xdata
                    cur_y = event.ydata
                    click = event.button
                    if click == 1:
                        print('Closing')
                        pyplot.close()
            connection_id = fig.canvas.mpl_connect('button_press_event', onclick)
            pyplot.savefig(folder_selected+"Figure_1.png")
            pyplot.show()

        deg_file = []
        sig_file = []
        if not auto:
            center_x = int(cur_x)
            center_y = int(cur_y)
            region_size = int(input('Enter the box size: '))
            half_region_size = (np.ceil(region_size / 2)).astype(int)
            SHG_Raw = np.loadtxt(folder_selected + file_name + "_{}deg".format(degree) + ".txt", dtype=int, delimiter=',')
            fig, ax = pyplot.subplots()
            region = SHG_Raw[center_x - half_region_size: center_x + half_region_size,
                                 center_y - half_region_size: center_y + half_region_size]
            im = ax.imshow(SHG_Raw, vmin=1000, vmax=5000)
            fig.colorbar(im, ax=ax, label='Interactive colorbar')
            ax.scatter(center_x, center_y, s=30, color='tomato', marker='x')
            rect = patches.Rectangle((center_x - half_region_size, center_y - half_region_size),
                                         region_size, region_size, linewidth=1, edgecolor='r', facecolor='none')
                # Add the patch to the Axes
            ax.add_patch(rect)
            pyplot.show()
            angle = 365
        for degree in range(0, 365, step_size):
            deg_file = np.append(deg_file, degree)
            SHG_Raw = np.loadtxt(folder_selected + file_name + "_{}deg".format(degree) + ".txt", dtype=int, delimiter=',')
            if not auto:
                region = SHG_Raw[center_x - half_region_size: center_x + half_region_size,
                                 center_y - half_region_size: center_y + half_region_size]
            else:
                region = SHG_Raw[top_left_col: top_left_col + region_size, top_left_row: top_left_row + region_size]

            small_sum = sum(map(sum, region))
            large_sum = sum(map(sum, SHG_Raw))
            bkg_pixel = (large_sum - small_sum) / (512 ** 2 - region_size ** 2)
            sig = small_sum - bkg_pixel * region_size ** 2
            sig_file = np.append(sig_file, sig)

        sig_file = sig_file.astype(np.float64)
        # sig_file = sig_file.tolist()
        max_lim = max(sig_file)
        min_lim = min(sig_file)
        deg_file = deg_file * np.pi / 180
        deg_file = deg_file.astype(np.float64)
        # deg_file = deg_file.tolist()
        fig, ax = pyplot.subplots(subplot_kw={'projection': 'polar'})
        ax.plot(deg_file, sig_file, color='red')
        ax.set_ylim(bottom=min_lim, top=max_lim)
        # pyplot.autoscale()
        pyplot.show()
        pyplot.plot(deg_file, sig_file, linewidth=5, color='blue')
        pyplot.close()
        csv_file_name = 'Processed_First_Data.csv'
        comb = pd.DataFrame(list(zip(deg_file, sig_file)))
        rec_data = pd.DataFrame()
        rec_data = pd.concat([rec_data, comb], ignore_index=True, axis=1)
        rec_data.to_csv(folder_selected + csv_file_name, mode='a', index=False, encoding='utf-8-sig', header=None)

        slope = (sig_file[-1] - sig_file[0]) / (deg_file[-1] - deg_file[0])
        const = sig_file[-1] - slope * deg_file[-1]
        # const = sig_file[-1] + slope * deg_file[-1]
        sig_file = sig_file - (slope * deg_file + const)
        sig_file = (sig_file/30)/380000
        csv_file_name = 'Processed_Data.csv'
        comb = pd.DataFrame(list(zip(deg_file, sig_file)))
        rec_data = pd.DataFrame()
        rec_data = pd.concat([rec_data, comb], ignore_index=True, axis=1)
        rec_data.to_csv(folder_selected + csv_file_name, mode='a', index=False, encoding='utf-8-sig', header=None)

        min_sig = min(sig_file)
        sig_file = sig_file - min_sig
        csv_file_name = 'Processed_Data_Min_Shift.csv'
        comb = pd.DataFrame(list(zip(deg_file, sig_file)))
        rec_data = pd.DataFrame()
        rec_data = pd.concat([rec_data, comb], ignore_index=True, axis=1)
        rec_data.to_csv(folder_selected + csv_file_name, mode='w', index=False, encoding='utf-8-sig', header=None)

        pyplot.plot(deg_file, sig_file, linewidth=5, color='blue')
        pyplot.show()

        max_lim = max(sig_file)
        min_lim = min(sig_file)


        fig, ax = pyplot.subplots(subplot_kw={'projection': 'polar'})
        ax.plot(deg_file, sig_file, color='tomato')
        ax.scatter(deg_file, sig_file, color='tomato')
        ax.set_ylim(bottom=min_lim, top=max_lim)
        pyplot.title(title + '{} Polarization'.format(polarization), pad=10, wrap=True)
        pyplot.tight_layout()
        pyplot.savefig(folder_selected+"Figure_2.png")
        pyplot.show()

        fit = True
        if fit == True:
            def shg_sin(params, x, data=None):
                A = params['A']
                a2 = params['a2']
                B = params['B']
                x0 = params['x0']
                model = (A*np.sin(a2-3*(x-x0))+B*np.sin(a2+(x-x0)))**2
                if data is None:
                    return model
                return model - data

            def shg_cos(params, x, data=None):


                A = params['A']
                a2 = params['a2']
                B = params['B']
                model = A*np.cos(a2-3*x)+B*np.cos(a2+3*x)*2
                if data is None:
                    return model
                return model - data

            # Create a Parameters object
            params = lmfit.Parameters()
            params.add('A', value=-0.1)
            params.add('a2', value=-0.1)
            params.add('B', value=11)
            params.add('x0', value=0.1)
            result_sin = lmfit.minimize(shg_sin, params, args=(deg_file,), kws={'data': sig_file})
            sin_A = result_sin.params['A'].value
            sin_a2 = result_sin.params['a2'].value
            sin_B = result_sin.params['B'].value
            sin_x0 = result_sin.params['x0'].value

            sin_A_err = result_sin.params['A'].stderr
            sin_a2_err = result_sin.params['a2'].stderr
            sin_B_err = result_sin.params['B'].stderr
            sin_x0_err = result_sin.params['x0'].stderr

            fig, ax = pyplot.subplots(subplot_kw={'projection': 'polar'})

            ax.plot(deg_file, result_sin.residual + sig_file, color='#E74C3C',linewidth=2)
            ax.scatter(deg_file, sig_file, color='black', s=4)
            ax.set_ylim(bottom=min_lim*1.1, top=max_lim*1.1)
            pyplot.title(title + '{} Polarization'.format(polarization), pad=10, wrap=True)
            pyplot.tight_layout()
            pyplot.savefig(folder_selected + "Fitted_Data.png")
            pyplot.show()

            df = pd.DataFrame()
            df_comb = pd.DataFrame(list(zip([sin_A], [sin_A_err], [sin_a2], [sin_a2_err], [sin_B], [sin_B_err],[sin_x0], [sin_x0_err])))
            df = pd.concat([df, df_comb], ignore_index=True, axis=1)
            df.to_csv(folder_selected + 'Fitted_Data.csv', index=False)

            # df = pd.DataFrame()
            # df_comb = pd.DataFrame(list(zip([sin_A_err], [sin_a2_err], [sin_B_err], [sin_x0_err])))
            # df = pd.concat([df, df_comb], ignore_index=True, axis=1)
            # df.to_csv(folder_selected + 'Fitted_Data.csv', index=False)

        if os.path.exists(folder_selected + 'Results.pptx'):
            prs = Presentation(folder_selected + 'Results.pptx')
            prs.slide_width = Inches(13.33)
            prs.slide_height = Inches(7.5)
        else:
            prs = Presentation()
            prs.slide_width = Inches(13.33)
            prs.slide_height = Inches(7.5)
            prs.save(folder_selected + 'Results.pptx')
            prs = Presentation(folder_selected + 'Results.pptx')
            prs.slide_width = Inches(13.33)
            prs.slide_height = Inches(7.5)

        SHG_Image = folder_selected + 'Figure_1.png'
        SHG_Signal = folder_selected + 'Figure_2.png'
        blank_slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_slide_layout)
        SHG_Image_img = slide.shapes.add_picture(SHG_Image, Inches(0.32), Inches(1.42), Inches(6.39))
        SHG_Signal_img = slide.shapes.add_picture(SHG_Signal, Inches(6.49), Inches(1.42), Inches(6.39))
        text_frame = slide.shapes.add_textbox(Inches(0.18), Inches(0.2), Inches(6.67), Inches(0.4))
        Data_frame = slide.shapes.add_textbox(Inches(11.19), Inches(0.2), Inches(2.04), Inches(0.4))
        text_frame = text_frame.text_frame
        Data_frame = Data_frame.text_frame
        p = text_frame.paragraphs[0]
        d = Data_frame.paragraphs[0]
        run = p.add_run()
        run.text = str(folder_name_pptx)
        font = run.font
        font.name = 'Calibri'
        font.size = Pt(18)

        run_d = d.add_run()
        run_d.text = str(Date)
        font_d = run_d.font
        font_d.name = 'Calibri'
        font_d.size = Pt(18)
        prs.save(folder_selected + 'Results.pptx')


if __name__ == '__main__':
    running = True
    while running:
    # construct the main wi
        try:
            window1 = SHG_Processing()
        except KeyboardInterrupt:
            break
    # endless loop unless quit
